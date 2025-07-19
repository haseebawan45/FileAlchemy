from fastapi import FastAPI, File, Query, UploadFile, APIRouter, BackgroundTasks, WebSocket
from fastapi.responses import FileResponse, JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
import os
import fitz  # PyMuPDF for PDF processing
import pandas as pd
from docx import Document
from docx.shared import Inches
import markdown
import xmltodict
import json
import subprocess
from weasyprint import HTML  # Faster HTML to PDF
import easyocr  # Faster OCR
from PIL import Image, ImageEnhance, ImageFilter, ImageDraw, ImageFont
from pptx import Presentation
# Conditionally import moviepy later
import aiofiles
import pdfplumber
import pdfkit
from pdf2image import convert_from_path
import pytesseract
import traceback
import numpy as np
from PIL import ImageStat
from skimage.filters import threshold_otsu
import asyncio
import time
import re
import io
import zipfile
import tempfile
import stat
import ctypes
import subprocess
import sys
import uuid
import shutil
import base64
import random
import string
import asyncio
import secrets
import hashlib
import traceback
import tempfile
import subprocess
from datetime import datetime, timedelta
from resource_manager import get_resource_manager
from optimized_converter import get_optimized_converter
from conversion_queue import get_conversion_queue, get_queue_processor, ConversionRequest

# Check if ffmpeg is available - needed for moviepy
ffmpeg_available = False
ffmpeg_path = None
try:
    import imageio_ffmpeg
    ffmpeg_path = imageio_ffmpeg.get_ffmpeg_exe()
    # Only import moviepy if ffmpeg is available
    from moviepy.video.io.ImageSequenceClip import ImageSequenceClip
    ffmpeg_available = True
    print(f"FFmpeg found at {ffmpeg_path} - video conversion features enabled")
except Exception as e:
    print(f"FFmpeg not found - attempting to locate local copy: {str(e)}")
    # Try to use local ffmpeg if available
    local_ffmpeg = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ffmpeg", "ffmpeg.exe")
    if os.path.exists(local_ffmpeg):
        print(f"Found local FFmpeg at {local_ffmpeg}")
        os.environ["IMAGEIO_FFMPEG_EXE"] = local_ffmpeg
        try:
            import imageio_ffmpeg
            ffmpeg_path = imageio_ffmpeg.get_ffmpeg_exe()
            from moviepy.video.io.ImageSequenceClip import ImageSequenceClip
            ffmpeg_available = True
            print(f"Successfully using local FFmpeg at {ffmpeg_path}")
        except Exception as e2:
            print(f"Failed to use local FFmpeg: {str(e2)}")
            # Create a dummy function to avoid errors when ffmpeg is not available
            class DummyImageSequenceClip:
                def __init__(self, *args, **kwargs):
                    raise RuntimeError("Video conversion is not available. FFmpeg is required but not found. Please install FFmpeg and add it to your PATH or set the IMAGEIO_FFMPEG_EXE environment variable to point to the ffmpeg executable.")
            ImageSequenceClip = DummyImageSequenceClip
    else:
        # Create a dummy function to avoid errors when ffmpeg is not available
        class DummyImageSequenceClip:
            def __init__(self, *args, **kwargs):
                raise RuntimeError("Video conversion is not available. FFmpeg is required but not found. Please install FFmpeg and add it to your PATH or set the IMAGEIO_FFMPEG_EXE environment variable to point to the ffmpeg executable.")
        ImageSequenceClip = DummyImageSequenceClip

# Check for lxml - needed for enhanced image extraction from PPTX
lxml_available = False
try:
    import lxml
    from lxml import etree
    lxml_available = True
    print("lxml found - enhanced PPTX image extraction enabled")
except Exception as e:
    print(f"lxml not found - some image extraction features will be limited: {str(e)}")
    # We'll handle this gracefully in the code

# Try to find ffmpeg in common locations if not found by imageio
if not ffmpeg_available:
    common_ffmpeg_locations = [
        # Windows
        r"C:\Program Files\ffmpeg\bin\ffmpeg.exe",
        r"C:\ffmpeg\bin\ffmpeg.exe",
        # Linux
        "/usr/bin/ffmpeg",
        "/usr/local/bin/ffmpeg",
        # macOS
        "/usr/local/bin/ffmpeg",
        "/opt/homebrew/bin/ffmpeg"
    ]
    
    for location in common_ffmpeg_locations:
        if os.path.exists(location):
            print(f"Found FFmpeg at {location}, trying to use it")
            os.environ["IMAGEIO_FFMPEG_EXE"] = location
            try:
                import imageio_ffmpeg
                from moviepy.video.io.ImageSequenceClip import ImageSequenceClip
                ffmpeg_available = True
                print("Successfully loaded FFmpeg and enabled video conversion features")
                break
            except Exception as e:
                print(f"Failed to use FFmpeg at {location}: {str(e)}")

# Define a simple alternative when FFmpeg is not available
def create_fallback_slideshow(image_paths, output_path):
    """Creates an animated GIF slideshow as fallback when FFmpeg is not available"""
    try:
        print("Creating GIF slideshow as fallback for video conversion")
        if not image_paths or len(image_paths) == 0:
            raise ValueError("No images provided for slideshow creation")
            
        # Load all images
        images = [Image.open(img_path) for img_path in image_paths]
        
        # Resize images for consistency
        max_size = (800, 600)
        for i in range(len(images)):
            images[i].thumbnail(max_size, Image.LANCZOS)
            
        # Save as animated GIF
        gif_path = output_path.replace(".mp4", ".gif")
        images[0].save(
            gif_path,
            save_all=True,
            append_images=images[1:],
            optimize=False,
            duration=1000,  # 1 second per frame
            loop=0  # Loop forever
        )
        print(f"Created GIF slideshow at {gif_path}")
        return gif_path
    except Exception as e:
        print(f"Failed to create fallback slideshow: {str(e)}")
        raise

app = FastAPI()

# Create an API router
api_router = APIRouter()

# Startup event to initialize queue processor
@app.on_event("startup")
async def startup_event():
    """Initialize services on startup"""
    # Start the queue processor
    asyncio.create_task(queue_processor.start_processing())
    print("Queue processor started")

@app.on_event("shutdown")
async def shutdown_event():
    """Cleanup on shutdown"""
    queue_processor.stop_processing()
    resource_manager.cleanup_resources()
    print("Services shut down")

# Railway environment configuration
RAILWAY_ENVIRONMENT = os.getenv('RAILWAY_ENVIRONMENT', 'development')
PORT = int(os.getenv('PORT', 8080))

# CORS Middleware - Configure for Railway deployment
allowed_origins = ["*"] if RAILWAY_ENVIRONMENT == 'development' else [
    "https://*.railway.app",
    "https://*.up.railway.app", 
    "https://haseebawan45.github.io"
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=allowed_origins,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Directories
UPLOAD_DIR = "uploads"
CONVERTED_DIR = "converted"

os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(CONVERTED_DIR, exist_ok=True)

# Initialize EasyOCR reader
ocr_reader = easyocr.Reader(['en'])

# Global dict to store conversion progress
conversion_progress = {}

# Get resource manager instance
resource_manager = get_resource_manager()

# Get optimized converter instance
optimized_converter = get_optimized_converter()

# Get queue instances
conversion_queue = get_conversion_queue()
queue_processor = get_queue_processor(resource_manager)

# Define API endpoints
@api_router.post("/convert/")
async def convert_file(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    target_format: str = Query(..., description="Target file format")
):
    try:
        print(f"Starting conversion: {file.filename} to {target_format}")
        
        # Generate a unique task ID for tracking progress
        task_id = f"conversion_{int(time.time())}_{os.urandom(4).hex()}"
        conversion_progress[task_id] = {"progress": 0, "status": "Starting conversion"}
        
        # Check if we can process this request
        estimated_memory = min(file_size // (1024 * 1024) * 2, 100)  # Estimate 2MB per MB of file, max 100MB
        if not resource_manager.can_process_request(estimated_memory):
            conversion_progress[task_id] = {
                "progress": 100, 
                "status": "Error: Server is at capacity. Please try again later."
            }
            return {"error": "Server is at capacity. Please try again later.", "task_id": task_id}
        
        # Check file size (limit to 25MB for Railway free tier)
        MAX_FILE_SIZE = 25 * 1024 * 1024  # 25MB in bytes
        
        # Read a small chunk to initialize file reading
        content = await file.read(1024)
        file_size = len(content)
        
        # Continue reading the file to determine total size
        while True:
            chunk = await file.read(1024)
            if not chunk:
                break
            file_size += len(chunk)
            
            # Check if exceeds maximum size
            if file_size > MAX_FILE_SIZE:
                print(f"File too large: {file_size} bytes")
                conversion_progress[task_id] = {"progress": 100, "status": "Error: File size exceeds maximum limit of 25MB"}
                return {"error": "File size exceeds maximum limit of 25MB", "task_id": task_id}
        
        # Reset the file position for subsequent reading
        await file.seek(0)
        
        # Update progress
        conversion_progress[task_id] = {"progress": 5, "status": "Saving uploaded file"}
        
        # Create paths if they don't exist
        os.makedirs(UPLOAD_DIR, exist_ok=True)
        os.makedirs(CONVERTED_DIR, exist_ok=True)
        
        file_path = os.path.join(UPLOAD_DIR, file.filename)
        print(f"Saving file to: {file_path}")
        
        # Save uploaded file asynchronously
        async with aiofiles.open(file_path, "wb") as buffer:
            content = await file.read()
            await buffer.write(content)
            print(f"File saved successfully, size: {len(content)} bytes")
        
        # Check if file was saved successfully
        conversion_progress[task_id] = {"progress": 10, "status": "File saved, starting conversion"}
        
        if not os.path.exists(file_path):
            print(f"Error: File not saved at {file_path}")
            conversion_progress[task_id] = {"progress": 100, "status": "Error: Failed to save uploaded file"}
            return {"error": "Failed to save uploaded file", "task_id": task_id}
            
        if os.path.getsize(file_path) == 0:
            print(f"Error: Saved file is empty at {file_path}")
            conversion_progress[task_id] = {"progress": 100, "status": "Error: Uploaded file is empty"}
            return {"error": "Uploaded file is empty", "task_id": task_id}
        
        extension = file.filename.split(".")[-1].lower()
        print(f"File extension detected: {extension}")
        
        # Prepare response with task_id for frontend progress tracking
        response_data = {"task_id": task_id}
        
        converted_file_path = None
        
        # Try optimized conversion methods first
        optimized_func = None
        if extension == "pdf":
            optimized_func = optimized_converter.select_conversion_method("pdf", target_format)
        elif extension in ["jpg", "jpeg", "png", "webp", "bmp", "tiff", "gif"]:
            if target_format in ["jpg", "jpeg", "png", "webp", "bmp", "tiff", "gif"]:
                optimized_func = lambda path, task_id=None: optimized_converter.streaming_image_convert(path, target_format, task_id)
        elif extension == "xlsx":
            optimized_func = optimized_converter.select_conversion_method("xlsx", target_format)
        
        if optimized_func:
            # Create conversion request for queue
            request = ConversionRequest(
                task_id=task_id,
                file_path=file_path,
                convert_func=optimized_func,
                estimated_memory_mb=estimated_memory,
                priority=1
            )
            
            # Try to add to queue
            if conversion_queue.add_request(request):
                conversion_progress[task_id] = {
                    "progress": 0,
                    "status": f"Queued for processing. Position: {conversion_queue.get_request_position(task_id)}"
                }
                return response_data
            else:
                conversion_progress[task_id] = {
                    "progress": 100,
                    "status": "Error: Server queue is full. Please try again later."
                }
                return {"error": "Server queue is full. Please try again later.", "task_id": task_id}
        
        # Fallback to original conversion functions
        elif extension in conversion_functions and target_format in conversion_functions[extension]:
            # Get the conversion function
            convert_func = conversion_functions[extension][target_format]
            
            # Create conversion request for queue
            request = ConversionRequest(
                task_id=task_id,
                file_path=file_path,
                convert_func=convert_func,
                estimated_memory_mb=estimated_memory,
                priority=1
            )
            
            # Try to add to queue
            if conversion_queue.add_request(request):
                conversion_progress[task_id] = {
                    "progress": 0,
                    "status": f"Queued for processing. Position: {conversion_queue.get_request_position(task_id)}"
                }
                return response_data
            else:
                conversion_progress[task_id] = {
                    "progress": 100,
                    "status": "Error: Server queue is full. Please try again later."
                }
                return {"error": "Server queue is full. Please try again later.", "task_id": task_id}
        # Special case for OCR - use optimized method
        elif target_format == "text (ocr)" and extension in ["jpg", "jpeg", "png", "webp", "bmp", "tiff"]:
            # Use optimized OCR method
            ocr_func = lambda path, task_id=None: optimized_converter.memory_efficient_ocr(path, task_id)
            
            # Create conversion request for queue
            request = ConversionRequest(
                task_id=task_id,
                file_path=file_path,
                convert_func=ocr_func,
                estimated_memory_mb=estimated_memory,
                priority=1
            )
            
            # Try to add to queue
            if conversion_queue.add_request(request):
                conversion_progress[task_id] = {
                    "progress": 0,
                    "status": f"Queued for processing. Position: {conversion_queue.get_request_position(task_id)}"
                }
                return response_data
            else:
                conversion_progress[task_id] = {
                    "progress": 100,
                    "status": "Error: Server queue is full. Please try again later."
                }
                return {"error": "Server queue is full. Please try again later.", "task_id": task_id}
        else:
            # Handle unsupported conversions
            conversion_progress[task_id] = {
                "progress": 100, 
                "status": f"Error: Conversion from {extension} to {target_format} is not supported"
            }
            return {
                "error": f"Conversion from {extension} to {target_format} is not supported",
                "task_id": task_id
            }
            
    except Exception as e:
        # Log the error 
        print(f"Conversion error: {str(e)}")
        task_id = task_id if 'task_id' in locals() else f"error_{int(time.time())}"
        conversion_progress[task_id] = {"progress": 100, "status": f"Error: {str(e)}"}
        return {"error": f"Conversion failed: {str(e)}", "task_id": task_id}

# Add this helper function for background processing
async def process_conversion(file_path, convert_func, task_id):
    """Process the conversion in a background task and update progress"""
    # Start conversion tracking
    if not resource_manager.start_conversion(task_id):
        conversion_progress[task_id] = {"progress": 100, "status": "Error: Server is at capacity"}
        return
    
    try:
        # Register temp file for cleanup
        resource_manager.register_temp_file(file_path)
        # Call the conversion function with task_id if supported
        import inspect
        sig = inspect.signature(convert_func)
        is_async = inspect.iscoroutinefunction(convert_func)
        
        if 'task_id' in sig.parameters:
            # Function supports progress tracking
            if is_async:
                converted_file_path = await convert_func(file_path, task_id=task_id)
            else:
                converted_file_path = convert_func(file_path, task_id=task_id)
        else:
            # Legacy function without progress tracking
            # We'll update progress manually
            conversion_progress[task_id] = {"progress": 20, "status": "Processing conversion"}
            if is_async:
                converted_file_path = await convert_func(file_path)
            else:
                if isinstance(convert_func, type(lambda: None)) and convert_func.__name__ == '<lambda>':
                    # Handle lambda functions that might return coroutines
                    result = convert_func(file_path)
                    if inspect.iscoroutine(result):
                        converted_file_path = await result
                    else:
                        converted_file_path = result
                else:
                    converted_file_path = convert_func(file_path)
            conversion_progress[task_id] = {"progress": 90, "status": "Conversion processing complete"}
        
        # Update final progress if not already updated by the conversion function
        if task_id in conversion_progress and converted_file_path:
            # Store the file path and name for download
            conversion_progress[task_id].update({
                "progress": 100, 
                "status": "Conversion complete",
                "file_path": converted_file_path,
                "file_name": os.path.basename(converted_file_path)
            })
        
        # Clean up the uploaded file
        if os.path.exists(file_path):
            os.remove(file_path)
            
    except Exception as e:
        print(f"Background conversion error: {str(e)}")
        traceback_info = traceback.format_exc()
        print(f"Traceback: {traceback_info}")
        conversion_progress[task_id] = {"progress": 100, "status": f"Error: {str(e)}"}
    finally:
        # Always end conversion tracking and cleanup
        resource_manager.end_conversion(task_id)

async def process_ocr(file_path, task_id):
    """Process OCR in a background task and update progress"""
    # Start conversion tracking
    if not resource_manager.start_conversion(task_id):
        conversion_progress[task_id] = {"progress": 100, "status": "Error: Server is at capacity"}
        return
    
    try:
        # Register temp file for cleanup
        resource_manager.register_temp_file(file_path)
        conversion_progress[task_id] = {"progress": 20, "status": "Processing OCR"}
        
        # Check if extract_text_from_image is async
        import inspect
        is_async = inspect.iscoroutinefunction(extract_text_from_image)
        
        if is_async:
            converted_file_path = await extract_text_from_image(file_path)
        else:
            converted_file_path = extract_text_from_image(file_path)
        
        # Update final progress
        conversion_progress[task_id] = {
            "progress": 100, 
            "status": "OCR complete",
            "file_path": converted_file_path,
            "file_name": os.path.basename(converted_file_path)
        }
        
        # Clean up the uploaded file
        if os.path.exists(file_path):
            os.remove(file_path)
            
    except Exception as e:
        print(f"OCR error: {str(e)}")
        conversion_progress[task_id] = {"progress": 100, "status": f"Error: {str(e)}"}
    finally:
        # Always end conversion tracking and cleanup
        resource_manager.end_conversion(task_id)

# Add a new endpoint to download converted files
@app.get("/download/{task_id}")
async def download_file(task_id: str):
    """Download a converted file by task ID"""
    if task_id in conversion_progress and "file_path" in conversion_progress[task_id]:
        progress_data = conversion_progress[task_id]
        file_path = progress_data["file_path"]
        
        if os.path.exists(file_path):
            # Check if it's a directory (for image conversions)
            if os.path.isdir(file_path):
                # Create a zip file of the directory
                # Create a temporary zip file
                zip_path = os.path.join(CONVERTED_DIR, f"{os.path.basename(file_path)}.zip")
                
                try:
                    with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                        # Walk through all files in the directory and add them to the zip
                        for root, _, files in os.walk(file_path):
                            for file in files:
                                file_full_path = os.path.join(root, file)
                                # Calculate the relative path for the zip structure
                                rel_path = os.path.relpath(file_full_path, file_path)
                                zipf.write(file_full_path, rel_path)
                    
                    # Return the zip file
                    return FileResponse(
                        zip_path,
                        filename=f"{os.path.basename(file_path)}.zip",
                        headers={"Access-Control-Expose-Headers": "Content-Disposition"}
                    )
                except Exception as e:
                    print(f"Error creating zip file: {str(e)}")
                    return JSONResponse(
                        status_code=500,
                        content={"error": f"Error creating zip file: {str(e)}"}
                    )
            else:
                # It's a regular file, return it directly
                return FileResponse(
                    file_path,
                    filename=progress_data.get("file_name", os.path.basename(file_path)),
                    headers={"Access-Control-Expose-Headers": "Content-Disposition"}
                )
    
    return JSONResponse(
        status_code=404,
        content={"error": "File not found or conversion not complete"}
    )

# Mount API router
app.include_router(api_router)

# Add endpoints for conversion progress tracking
@app.get("/conversion-progress/{task_id}")
async def get_conversion_progress(task_id: str):
    """Get the current progress of a conversion task"""
    if task_id in conversion_progress:
        return conversion_progress[task_id]
    return {"progress": 0, "status": "Task not found"}

# Resource status endpoint
@app.get("/resource-status")
async def get_resource_status():
    """Get current resource usage status"""
    return resource_manager.get_status()

# Queue status endpoint
@app.get("/queue-status")
async def get_queue_status():
    """Get current queue status"""
    return conversion_queue.get_queue_status()

# Request position endpoint
@app.get("/queue-position/{task_id}")
async def get_queue_position(task_id: str):
    """Get the position of a request in the queue"""
    position = conversion_queue.get_request_position(task_id)
    if position:
        return {"task_id": task_id, "position": position}
    return {"task_id": task_id, "position": None, "status": "not_in_queue"}

# Health check endpoint for Railway
@app.get("/health")
async def health_check():
    return {"status": "healthy", "environment": RAILWAY_ENVIRONMENT}

# Root path handler for index.html
@app.get("/")
async def read_index():
    return FileResponse("index.html")

# Mount static files for specific directories
app.mount("/css", StaticFiles(directory="css"), name="css")
app.mount("/js", StaticFiles(directory="js"), name="js")
app.mount("/assets", StaticFiles(directory="assets"), name="assets")
app.mount("/pages", StaticFiles(directory="pages"), name="pages")

# Add route for firebase-config.js
@app.get("/firebase-config.js")
async def get_firebase_config():
    return FileResponse("firebase-config.js", media_type="application/javascript")

# =================== FILE CONVERSION FUNCTIONS =====================

# Check for some common issues in PDF conversions
def check_pdf_issues(pdf_path):
    if not os.path.exists(pdf_path):
        print(f"PDF file does not exist: {pdf_path}")
        return "PDF file does not exist"
        
    if os.path.getsize(pdf_path) == 0:
        print(f"PDF file is empty: {pdf_path}")
        return "PDF file is empty"
        
    try:
        # Try to open with PyMuPDF to check if it's valid
        pdf_document = fitz.open(pdf_path)
        num_pages = len(pdf_document)
        print(f"PDF validation successful: {num_pages} pages")
        pdf_document.close()
        return None
    except Exception as e:
        print(f"PDF validation failed: {str(e)}")
        return f"Invalid PDF file: {str(e)}"

# ✅ PDF to DOCX (Using PyMuPDF)
async def convert_pdf_to_docx(pdf_path):
    try:
        print(f"Starting PDF to DOCX conversion: {pdf_path}")
        
        # Check for common PDF issues
        pdf_issue = check_pdf_issues(pdf_path)
        if pdf_issue:
            raise Exception(pdf_issue)
        
        docx_path = os.path.join(CONVERTED_DIR, os.path.basename(pdf_path).replace(".pdf", "_converted.docx"))
        print(f"Output will be saved to: {docx_path}")
        
        # First, try to use pdf2docx if available (better formatting preservation)
        try:
            from pdf2docx import Converter
            print("Using pdf2docx for better formatting preservation")
            
            # Convert with pdf2docx
            cv = Converter(pdf_path)
            cv.convert(docx_path, start=0, end=None)
            cv.close()
            
            # Check if conversion was successful
            if os.path.exists(docx_path) and os.path.getsize(docx_path) > 0:
                print(f"PDF to DOCX conversion completed successfully using pdf2docx: {docx_path}")
                return docx_path
            else:
                print("pdf2docx conversion failed or produced empty file, falling back to PyMuPDF")
        except Exception as e:
            print(f"pdf2docx conversion failed, falling back to PyMuPDF: {str(e)}")
        
        # Fallback to PyMuPDF approach
        print("Using PyMuPDF for PDF to DOCX conversion")
        
        # Open PDF document
        pdf_document = fitz.open(pdf_path)
        doc = Document()
        
        # Add a title 
        doc.add_heading(os.path.basename(pdf_path).replace(".pdf", ""), 0)
        
        # Create a temporary directory for extracted images
        temp_image_dir = os.path.join(CONVERTED_DIR, "temp_images")
        os.makedirs(temp_image_dir, exist_ok=True)
        image_count = 0
        
        # Process each page
        for page_num in range(len(pdf_document)):
            print(f"Processing page {page_num+1}/{len(pdf_document)}")
            page = pdf_document.load_page(page_num)
            
            # Extract page dimensions for layout analysis
            page_width = page.rect.width
            page_height = page.rect.height
            
            # Add a page heading
            if page_num > 0:  # Not for the first page since we already added document title
                doc.add_heading(f"Page {page_num + 1}", level=1)
            
            # Extract text blocks with more detailed structure info
            blocks = page.get_text("dict")["blocks"]
            print(f"Extracted {len(blocks)} content blocks from page {page_num+1}")
            
            # Process blocks by type
            for block in blocks:
                # Image blocks
                if block["type"] == 1:  # Image
                    try:
                        # Extract image
                        xref = block.get("xref", 0)
                        if xref > 0:  # Valid image reference
                            image_bytes = pdf_document.extract_image(xref)
                            if image_bytes:
                                image_data = image_bytes["image"]
                                image_ext = image_bytes["ext"]
                                
                                # Save image to temp dir
                                image_path = os.path.join(temp_image_dir, f"image_{image_count}.{image_ext}")
                                with open(image_path, "wb") as img_file:
                                    img_file.write(image_data)
                                
                                # Add image to document
                                doc.add_picture(image_path, width=Inches(6))  # Adjust width as needed
                                image_count += 1
                    except Exception as img_err:
                        print(f"Warning: Failed to extract image: {str(img_err)}")
                
                # Text blocks
                elif block["type"] == 0:  # Text
                    if "lines" in block:
                        # Process text by lines
                        for line in block["lines"]:
                            if "spans" in line:
                                line_text = ""
                                is_bold = False
                                is_italic = False
                                font_size = 0
                                
                                # Process text spans to preserve formatting
                                for span in line["spans"]:
                                    span_text = span.get("text", "").strip()
                                    if not span_text:
                                        continue
                                        
                                    # Check for formatting
                                    current_font = span.get("font", "").lower()
                                    current_size = span.get("size", 0)
                                    
                                    # Detect formatting based on font name
                                    current_bold = "bold" in current_font or "heavy" in current_font
                                    current_italic = "italic" in current_font or "oblique" in current_font
                                    
                                    # If formatting changes, add current text and start new format
                                    if (is_bold != current_bold or 
                                        is_italic != current_italic or 
                                        abs(font_size - current_size) > 1):  # Font size changed
                                        
                                        # Add existing text if any
                                        if line_text:
                                            p = doc.add_paragraph()
                                            run = p.add_run(line_text)
                                            run.bold = is_bold
                                            run.italic = is_italic
                                            line_text = ""
                                        
                                        # Update formatting
                                        is_bold = current_bold
                                        is_italic = current_italic
                                        font_size = current_size
                                    
                                    # Add current span text
                                    if line_text:
                                        line_text += " " + span_text
                                    else:
                                        line_text = span_text
                                
                                # Add remaining text in the line
                                if line_text:
                                    p = doc.add_paragraph()
                                    run = p.add_run(line_text)
                                    run.bold = is_bold
                                    run.italic = is_italic
                
                # Table blocks (for future enhancement)
                # We could detect tables based on layout analysis
            
            # Add a page break after each page except the last one
            if page_num < len(pdf_document) - 1:
                doc.add_page_break()
        
        # Clean up temp images
        try:
            for img_file in os.listdir(temp_image_dir):
                os.remove(os.path.join(temp_image_dir, img_file))
            os.rmdir(temp_image_dir)
        except Exception as e:
            print(f"Warning: Failed to clean up temp images: {str(e)}")
        
        # Save the document
        print(f"Saving DOCX file to: {docx_path}")
        doc.save(docx_path)
        
        # Verify the file was created properly
        if not os.path.exists(docx_path):
            raise Exception(f"DOCX file was not created at {docx_path}")
            
        if os.path.getsize(docx_path) == 0:
            raise Exception(f"Generated DOCX file is empty: {docx_path}")
            
        print(f"PDF to DOCX conversion completed successfully: {docx_path}")
        return docx_path
        
    except Exception as e:
        print(f"Error in PDF to DOCX conversion: {str(e)}")
        traceback_info = traceback.format_exc()
        print(f"Traceback: {traceback_info}")
        raise Exception(f"PDF to DOCX conversion failed: {str(e)}")

# ✅ PDF to Excel (Extract Tables)
async def convert_pdf_to_excel(pdf_path):
    xlsx_path = os.path.join(CONVERTED_DIR, os.path.basename(pdf_path).replace(".pdf", "_converted.xlsx"))
    with pdfplumber.open(pdf_path) as pdf:
        tables = []
        for page in pdf.pages:
            tables.extend(page.extract_tables())
        df = pd.DataFrame(tables[0]) if tables else pd.DataFrame()
        df.to_excel(xlsx_path, index=False)
    return xlsx_path

# ✅ PDF to PowerPoint
async def convert_pdf_to_pptx(pdf_path):
    pptx_path = os.path.join(CONVERTED_DIR, os.path.basename(pdf_path).replace(".pdf", "_converted.pptx"))
    prs = Presentation()
    images = convert_from_path(pdf_path)
    for img in images:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        img_path = os.path.join(CONVERTED_DIR, "temp.png")
        img.save(img_path)
        slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)
    prs.save(pptx_path)
    return pptx_path

# ✅ PDF to Text (Using PyMuPDF)
async def convert_pdf_to_text(pdf_path):
    txt_path = os.path.join(CONVERTED_DIR, os.path.basename(pdf_path).replace(".pdf", "_converted.txt"))
    pdf_document = fitz.open(pdf_path)
    text = "\n".join([page.get_text("text") for page in pdf_document])
    
    async with aiofiles.open(txt_path, "w", encoding="utf-8") as f:
        await f.write(text)
    return txt_path

# ✅ PDF to HTML (Using pdfkit)
async def convert_pdf_to_html(pdf_path):
    try:
        print(f"Starting PDF to HTML conversion: {pdf_path}")
        html_path = os.path.join(CONVERTED_DIR, os.path.basename(pdf_path).replace(".pdf", "_converted.html"))
        print(f"Output will be saved to: {html_path}")
        
        try:
            # Try using pdfkit first
            print("Attempting conversion with pdfkit...")
            pdfkit.from_file(pdf_path, html_path)
            
            if os.path.exists(html_path) and os.path.getsize(html_path) > 0:
                print("PDF to HTML conversion with pdfkit successful")
                return html_path
            else:
                print("Pdfkit produced empty output file, trying alternative method...")
        except Exception as e:
            print(f"Pdfkit conversion failed: {str(e)}, trying alternative method...")
        
        # Alternative method using PyMuPDF if pdfkit fails
        print("Using PyMuPDF for PDF to HTML conversion...")
        pdf_document = fitz.open(pdf_path)
        html_content = ["<!DOCTYPE html><html><head><title>Converted PDF</title>",
                        "<style>body{font-family: Arial, sans-serif; margin: 40px;}</style></head><body>"]
        
        for page_num in range(len(pdf_document)):
            print(f"Processing page {page_num+1}/{len(pdf_document)}")
            page = pdf_document.load_page(page_num)
            html_content.append(f"<div class='page' id='page-{page_num+1}'>")
            html_content.append(f"<h2>Page {page_num+1}</h2>")
            
            # Get text as HTML
            text_html = page.get_text("html")
            html_content.append(text_html)
            html_content.append("</div><hr>")
        
        html_content.append("</body></html>")
        
        # Write the HTML content to file
        async with aiofiles.open(html_path, "w", encoding="utf-8") as f:
            await f.write("\n".join(html_content))
        
        print(f"PDF to HTML conversion with PyMuPDF successful")
        return html_path
        
    except Exception as e:
        print(f"Error in PDF to HTML conversion: {str(e)}")
        traceback_info = traceback.format_exc()
        print(f"Traceback: {traceback_info}")
        raise Exception(f"PDF to HTML conversion failed: {str(e)}")

# ✅ Image to OCR (Using EasyOCR)
async def extract_text_from_image(image_path, languages=['en']):
    try:
        print(f"Starting OCR text extraction from image: {image_path}")
        
        # Check if image exists
        if not os.path.exists(image_path):
            error_msg = f"Image file does not exist: {image_path}"
            print(error_msg)
            raise Exception(error_msg)
            
        # Check if image is not empty
        if os.path.getsize(image_path) == 0:
            error_msg = f"Image file is empty: {image_path}"
            print(error_msg)
            raise Exception(error_msg)
            
        txt_path = os.path.join(CONVERTED_DIR, os.path.splitext(os.path.basename(image_path))[0] + "_converted.txt")
        print(f"Text output will be saved to: {txt_path}")
        
        # Open and preprocess the image for better OCR results
        try:
            image = Image.open(image_path)
            print(f"Image opened successfully: format={image.format}, size={image.size}, mode={image.mode}")
        except Exception as e:
            error_msg = f"Failed to open image: {str(e)}"
            print(error_msg)
            raise Exception(error_msg)
        
        # Get basic image stats to determine preprocessing approach
        img_array = np.array(image)
        if len(img_array.shape) > 2:  # Color image
            # Calculate image contrast
            img_gray = image.convert('L')
            img_stats = ImageStat.Stat(img_gray)
            contrast = img_stats.stddev[0]
            
            # Calculate image brightness
            brightness = img_stats.mean[0]
            
            print(f"Image stats - contrast: {contrast:.2f}, brightness: {brightness:.2f}")
            
            # Apply appropriate preprocessing based on image properties
            if brightness < 100:  # Dark image
                print("Applying brightness enhancement for dark image")
                enhancer = ImageEnhance.Brightness(image)
                image = enhancer.enhance(1.5)
            
            if contrast < 40:  # Low contrast image
                print("Applying contrast enhancement for low contrast image")
                enhancer = ImageEnhance.Contrast(image)
                image = enhancer.enhance(1.5)
        
        # Convert to grayscale to improve OCR accuracy
        if image.mode != 'L':
            try:
                image = image.convert('L')
                print(f"Converted image to grayscale")
            except Exception as e:
                print(f"Warning: Failed to convert to grayscale: {str(e)}")
                # Continue with original image
            
        # Resize very large images to speed up processing
        max_dimension = 2000  # Max dimension in pixels
        if max(image.size) > max_dimension:
            try:
                ratio = max_dimension / max(image.size)
                new_size = (int(image.size[0] * ratio), int(image.size[1] * ratio))
                image = image.resize(new_size, Image.LANCZOS)
                print(f"Resized image to {new_size} for faster processing")
            except Exception as e:
                print(f"Warning: Failed to resize image: {str(e)}")
                # Continue with original image

        # Apply noise reduction and sharpening
        try:
            # Apply median filter to reduce noise
            image = image.filter(ImageFilter.MedianFilter(size=3))
            
            # Apply slight sharpening to enhance text edges
            image = image.filter(ImageFilter.SHARPEN)
            print("Applied noise reduction and sharpening")
        except Exception as e:
            print(f"Warning: Failed to apply filters: {str(e)}")

        # Apply thresholding for better text contrast
        try:
            # Use adaptive thresholding
            image_array = np.array(image)
            if image_array.ndim == 2:  # Ensure we have a grayscale image
                # Calculate optimal threshold using Otsu's method
                threshold_value = threshold_otsu(image_array)
                binary_image = (image_array > threshold_value) * 255
                image = Image.fromarray(binary_image.astype(np.uint8))
                print(f"Applied adaptive thresholding with value {threshold_value}")
        except Exception as e:
            print(f"Warning: Failed to apply thresholding: {str(e)}")

        # Save preprocessed image for debugging (optional)
        debug_image_path = os.path.join(CONVERTED_DIR, os.path.splitext(os.path.basename(image_path))[0] + "_preprocessed.png")
        image.save(debug_image_path)
        print(f"Saved preprocessed image to {debug_image_path} for reference")
            
        # Initialize EasyOCR reader with specified languages
        print(f"Initializing EasyOCR with languages: {languages}")
        reader = easyocr.Reader(languages, gpu=False)
        
        # Detect text with confidence scores
        results = reader.readtext(np.array(image))
        print(f"OCR detected {len(results)} text regions")
        
        # Get ordered text with confidence scores
        extracted_text = []
        for (bbox, text, confidence) in results:
            if confidence > 0.2:  # Filter low-confidence results
                extracted_text.append(f"{text} (confidence: {confidence:.2f})")
                print(f"Detected text: '{text}' with confidence {confidence:.2f}")
            
        # If EasyOCR didn't work well, try Tesseract as fallback
        if not extracted_text or len(extracted_text) < 3:
            print("EasyOCR results limited. Trying Tesseract as fallback...")
            try:
                tesseract_text = pytesseract.image_to_string(image)
                if tesseract_text.strip():
                    extracted_text = [tesseract_text + " (via Tesseract)"]
                    print("Successfully extracted text with Tesseract")
            except Exception as e:
                print(f"Tesseract fallback failed: {str(e)}")
        
        # Write final output
        final_text = "\n\n".join(extracted_text) if extracted_text else "No text detected in image"
        
        # Write to text file
        with open(txt_path, "w", encoding="utf-8") as f:
            f.write(final_text)
            
        print(f"Successfully wrote extracted text to {txt_path}")
        return txt_path
    
    except Exception as e:
        print(f"Error in image OCR processing: {str(e)}")
        traceback_info = traceback.format_exc()
        print(f"Traceback: {traceback_info}")
        raise Exception(f"OCR text extraction failed: {str(e)}")

# ✅ JSON to XML
def convert_json_to_xml(json_path):
    xml_path = os.path.join(CONVERTED_DIR, os.path.basename(json_path).replace(".json", ".xml"))
    with open(json_path, "r", encoding="utf-8") as f:
        json_data = json.load(f)
    xml_data = xmltodict.unparse({"root": json_data}, pretty=True)
    with open(xml_path, "w", encoding="utf-8") as f:
        f.write(xml_data)
    return xml_path

# ✅ XML to JSON
def convert_xml_to_json(xml_path):
    json_path = os.path.join(CONVERTED_DIR, os.path.basename(xml_path).replace(".xml", ".json"))
    with open(xml_path, "r", encoding="utf-8") as f:
        xml_data = f.read()
    json_data = json.dumps(xmltodict.parse(xml_data), indent=4)
    with open(json_path, "w", encoding="utf-8") as f:
        f.write(json_data)
    return json_path


# ✅ DOCX to PDF
def convert_docx_to_pdf(docx_path):
    try:
        print(f"Starting DOCX to PDF conversion: {docx_path}")
        pdf_path = os.path.join(CONVERTED_DIR, os.path.basename(docx_path).replace(".docx", ".pdf"))
        print(f"Output will be saved to: {pdf_path}")
        
        # First try using docx2pdf if available
        try:
            from docx2pdf import convert
            print("Attempting conversion with docx2pdf...")
            convert(docx_path, pdf_path)
            
            if os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 0:
                print(f"PDF created successfully with docx2pdf: {pdf_path}")
                return pdf_path
        except Exception as e:
            print(f"docx2pdf conversion failed: {str(e)}")
            print("Trying alternative method...")
        
        # Try with WeasyPrint
        try:
            print("Attempting conversion with WeasyPrint...")
            # First convert DOCX to HTML with our improved converter
            html_path = convert_docx_to_html(docx_path)
            
            if os.path.exists(html_path) and os.path.getsize(html_path) > 0:
                # Convert HTML to PDF using WeasyPrint
                HTML(html_path).write_pdf(pdf_path)
                
                # Clean up temporary HTML file
                os.remove(html_path)
                
                if os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 0:
                    print(f"PDF created successfully with WeasyPrint: {pdf_path}")
                    return pdf_path
            
            print("WeasyPrint conversion failed or produced empty file")
        except Exception as e:
            print(f"WeasyPrint conversion failed: {str(e)}")
            print("Trying next method...")
        
        # Try with pdfkit/wkhtmltopdf
        try:
            print("Attempting conversion with pdfkit/wkhtmltopdf...")
            # Convert DOCX to HTML if not already done
            if not os.path.exists(html_path):
                html_path = convert_docx_to_html(docx_path)
            
            # Configure pdfkit options for better output
            options = {
                'page-size': 'A4',
                'margin-top': '0.75in',
                'margin-right': '0.75in',
                'margin-bottom': '0.75in',
                'margin-left': '0.75in',
                'encoding': 'UTF-8',
                'no-outline': None
            }
            
            pdfkit.from_file(html_path, pdf_path, options=options)
            
            # Clean up temporary HTML file
            if os.path.exists(html_path):
                os.remove(html_path)
                
            if os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 0:
                print(f"PDF created successfully with pdfkit: {pdf_path}")
                return pdf_path
                
        except Exception as e:
            print(f"pdfkit conversion failed: {str(e)}")
            print("Trying final method...")
        
        # Final attempt with ReportLab
        try:
            print("Using ReportLab as final method...")
            
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import inch
            
            # Read DOCX file
            doc = Document(docx_path)
            
            # Create PDF
            document = SimpleDocTemplate(
                pdf_path,
                pagesize=letter,
                rightMargin=72,
                leftMargin=72,
                topMargin=72,
                bottomMargin=72
            )
            
            # Create styles
            styles = getSampleStyleSheet()
            styles.add(ParagraphStyle(
                name='CustomBody',
                parent=styles['Normal'],
                spaceBefore=6,
                spaceAfter=6,
                leading=14
            ))
            
            # Process document content
            elements = []
            
            for para in doc.paragraphs:
                if not para.text.strip():
                    continue
                    
                try:
                    # Clean text of any problematic characters
                    text = para.text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                    
                    # Determine style
                    if para.style.name.startswith('Heading'):
                        level = int(para.style.name[-1]) if para.style.name[-1].isdigit() else 1
                        style = styles[f'Heading{min(level, 4)}']
                    else:
                        style = styles['CustomBody']
                    
                    # Add paragraph
                    elements.append(Paragraph(text, style))
                    
                except Exception as para_error:
                    print(f"Warning: Skipping problematic paragraph: {str(para_error)}")
                    continue
            
            # Build the PDF
            document.build(elements)
            
            if os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 0:
                print(f"PDF created successfully with ReportLab: {pdf_path}")
                return pdf_path
            
        except Exception as e:
            print(f"ReportLab conversion failed: {str(e)}")
            raise Exception("All PDF conversion methods failed")
    
    except Exception as e:
        print(f"Error in DOCX to PDF conversion: {str(e)}")
        traceback_info = traceback.format_exc()
        print(f"Traceback: {traceback_info}")
        raise Exception(f"DOCX to PDF conversion failed: {str(e)}")

# ✅ DOCX to HTML
def convert_docx_to_html(docx_path):
    try:
        import base64
        from docx.shared import RGBColor, Pt, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.enum.style import WD_STYLE_TYPE
        from io import BytesIO

        html_path = os.path.join(CONVERTED_DIR, os.path.basename(docx_path).replace(".docx", ".html"))
        doc = Document(docx_path)
        
        # Enhanced HTML structure with complete styling
        html_content = [
            "<!DOCTYPE html>",
            "<html>",
            "<head>",
            "<meta charset='utf-8'>",
            "<style>",
            # Base styles with precise measurements
            """
            body {
                font-family: Calibri, Arial, sans-serif;
                line-height: 1.15;
                margin: 0 auto;
                padding: 20px;
                max-width: 8.5in;
                background: white;
            }
            .document {
                position: relative;
                width: 8.5in;
                margin: 0 auto;
                box-sizing: border-box;
                word-wrap: break-word;
            }
            p {
                margin: 0;
                padding: 0;
                line-height: inherit;
            }
            table {
                border-collapse: collapse;
                width: 100%;
            }
            td, th {
                vertical-align: top;
                padding: 0;
            }
            img {
                max-width: 100%;
                height: auto;
            }
            .section {
                position: relative;
                margin: 0;
                padding: 0;
            }
            .page-break {
                page-break-after: always;
            }
            """,
            "</style>",
            "</head>",
            "<body><div class='document'>"
        ]

        def rgb_to_hex(rgb):
            """Convert RGB color to hex format"""
            try:
                if isinstance(rgb, RGBColor):
                    # Get the RGB value as an integer
                    rgb_value = int(rgb._ms_color) if hasattr(rgb, '_ms_color') else None
                    if rgb_value is not None:
                        # Extract RGB components
                        r = (rgb_value >> 16) & 255
                        g = (rgb_value >> 8) & 255
                        b = rgb_value & 255
                        return f"#{r:02x}{g:02x}{b:02x}"
                elif isinstance(rgb, int):
                    # If it's already an integer value
                    r = (rgb >> 16) & 255
                    g = (rgb >> 8) & 255
                    b = rgb & 255
                    return f"#{r:02x}{g:02x}{b:02x}"
                return None
            except Exception as e:
                print(f"Warning: Color conversion failed: {str(e)}")
                return None

        def get_font_style(run):
            """Extract complete font styling from a run"""
            styles = []
            
            try:
                # Font family
                if hasattr(run, 'font') and run.font.name:
                    styles.append(f"font-family: '{run.font.name}'")

                # Font size in points
                if hasattr(run.font, 'size') and run.font.size:
                    size = run.font.size
                    if isinstance(size, Pt):
                        styles.append(f"font-size: {size.pt}pt")
                    else:
                        styles.append(f"font-size: {size/12700}pt")  # Convert from twips

                # Font color with improved handling
                if hasattr(run.font, 'color') and run.font.color and run.font.color.rgb:
                    color = rgb_to_hex(run.font.color.rgb)
                    if color:
                        styles.append(f"color: {color}")

                # Rest of the styling code...
                # Font color
                if hasattr(run.font, 'color') and run.font.color and run.font.color.rgb:
                    color = rgb_to_hex(run.font.color.rgb)
                    if color:
                        styles.append(f"color: {color}")

                # Text decorations
                if run.bold:
                    styles.append("font-weight: bold")
                if run.italic:
                    styles.append("font-style: italic")
                if run.underline:
                    styles.append("text-decoration: underline")
                if hasattr(run.font, 'strike') and run.font.strike:
                    styles.append("text-decoration: line-through")
                
                # Text effects
                if run.font.highlight_color:
                    styles.append(f"background-color: {run.font.highlight_color.lower()}")
                
                # Character spacing
                if hasattr(run.font, 'spacing'):
                    spacing = run.font.spacing
                    if spacing:
                        styles.append(f"letter-spacing: {spacing/20}pt")

                return "; ".join(styles)
            except Exception as e:
                print(f"Warning: Style extraction failed: {str(e)}")
                return ""

        def process_paragraph(paragraph):
            """Process paragraph with precise formatting"""
            if not paragraph.text.strip() and not paragraph.runs:
                return "<p style='margin-bottom: 1em'>&nbsp;</p>"

            # Get paragraph formatting
            p_format = paragraph.paragraph_format
            style_props = []

            # Alignment
            align_map = {
                WD_ALIGN_PARAGRAPH.LEFT: "left",
                WD_ALIGN_PARAGRAPH.CENTER: "center",
                WD_ALIGN_PARAGRAPH.RIGHT: "right",
                WD_ALIGN_PARAGRAPH.JUSTIFY: "justify"
            }
            if paragraph.alignment in align_map:
                style_props.append(f"text-align: {align_map[paragraph.alignment]}")

            # Indentation
            if p_format.first_line_indent:
                indent = p_format.first_line_indent
                if isinstance(indent, Inches):
                    style_props.append(f"text-indent: {indent.inches}in")
                else:
                    style_props.append(f"text-indent: {indent/914400}in")

            # Spacing
            if p_format.space_before:
                space = p_format.space_before
                if isinstance(space, Pt):
                    style_props.append(f"margin-top: {space.pt}pt")
                else:
                    style_props.append(f"margin-top: {space/20}pt")

            if p_format.space_after:
                space = p_format.space_after
                if isinstance(space, Pt):
                    style_props.append(f"margin-bottom: {space.pt}pt")
                else:
                    style_props.append(f"margin-bottom: {space/20}pt")

            if p_format.line_spacing:
                style_props.append(f"line-height: {p_format.line_spacing}")

            # Create paragraph HTML
            style_attr = f" style='{'; '.join(style_props)}'" if style_props else ""
            html = [f"<p{style_attr}>"]

            # Process runs with precise formatting
            for run in paragraph.runs:
                text = run.text
                if not text:
                    continue

                # Handle special characters
                text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')

                # Process images in runs
                if run.element.findall(".//w:drawing", namespaces=run.element.nsmap):
                    for inline_shape in paragraph.part.inline_shapes:
                        if inline_shape.width and inline_shape.height:
                            # Get image data and convert to base64
                            try:
                                image_data = inline_shape.blob
                                if image_data:
                                    img_base64 = base64.b64encode(image_data).decode()
                                    width = inline_shape.width / 914400  # Convert to inches
                                    height = inline_shape.height / 914400
                                    html.append(
                                        f"<img src='data:image/png;base64,{img_base64}' "
                                        f"style='width: {width}in; height: {height}in; "
                                        f"display: inline-block; vertical-align: text-bottom'>"
                                    )
                            except Exception as e:
                                print(f"Warning: Failed to process image: {str(e)}")
                    continue

                # Add run-level formatting
                style = get_font_style(run)
                if style:
                    text = f"<span style='{style}'>{text}</span>"

                html.append(text)

            html.append("</p>")
            return "".join(html)

        def process_table(table):
            """Process table with complete styling"""
            table_props = []
            
            # Table width
            if hasattr(table, 'width'):
                width = table.width
                if isinstance(width, Inches):
                    table_props.append(f"width: {width.inches}in")

            # Table borders
            borders = {'top': 0, 'right': 0, 'bottom': 0, 'left': 0}
            if hasattr(table, 'style'):
                for edge, border in borders.items():
                    if hasattr(table.style, f'{edge}_border'):
                        border_obj = getattr(table.style, f'{edge}_border')
                        if border_obj and border_obj.width:
                            borders[edge] = border_obj.width

            if any(borders.values()):
                table_props.append(f"border-width: {borders['top']}pt {borders['right']}pt "
                                 f"{borders['bottom']}pt {borders['left']}pt")

            style_attr = f" style='{'; '.join(table_props)}'" if table_props else ""
            html = [f"<table{style_attr}>"]

            # Process rows
            for row in table.rows:
                html.append("<tr>")
                for cell in row.cells:
                    # Cell properties
                    cell_props = []
                    
                    # Width and height
                    if cell.width:
                        cell_props.append(f"width: {cell.width.inches}in")
                    if hasattr(cell, 'height'):
                        cell_props.append(f"height: {cell.height.inches}in")

                    # Borders
                    cell_borders = []
                    for edge in ['top', 'right', 'bottom', 'left']:
                        border = getattr(cell, f'{edge}_border', None)
                        if border and border.width:
                            cell_borders.append(
                                f"border-{edge}: {border.width}pt solid {rgb_to_hex(border.color) or '#000'}"
                            )
                    if cell_borders:
                        cell_props.extend(cell_borders)

                    # Background color
                    if cell.fill.background_color:
                        bg_color = rgb_to_hex(cell.fill.background_color)
                        if bg_color:
                            cell_props.append(f"background-color: {bg_color}")

                    # Vertical alignment
                    if cell.vertical_alignment:
                        cell_props.append(f"vertical-align: {cell.vertical_alignment.lower()}")

                    # Add cell with styling
                    cell_style = f" style='{'; '.join(cell_props)}'" if cell_props else ""
                    html.append(f"<td{cell_style}>")

                    # Process cell content
                    for paragraph in cell.paragraphs:
                        html.append(process_paragraph(paragraph))

                    html.append("</td>")
                html.append("</tr>")

            html.append("</table>")
            return "\n".join(html)

        # Process document content
        for element in doc.element.body:
            if element.tag.endswith('}p'):
                para = [p for p in doc.paragraphs if p._element == element][0]
                html_content.append(process_paragraph(para))
            elif element.tag.endswith('}tbl'):
                table = [t for t in doc.tables if t._element == element][0]
                html_content.append(process_table(table))
            elif element.tag.endswith('}sectPr'):
                html_content.append("<div class='page-break'></div>")

        html_content.extend(["</div></body>", "</html>"])

        # Write the HTML file
        with open(html_path, "w", encoding="utf-8") as f:
            f.write("\n".join(html_content))

        return html_path

    except Exception as e:
        print(f"Error in DOCX to HTML conversion: {str(e)}")
        traceback_info = traceback.format_exc()
        print(f"Traceback: {traceback_info}")
        raise Exception(f"DOCX to HTML conversion failed: {str(e)}")

def convert_docx_to_markdown(docx_path):
    try:
        md_path = os.path.join(CONVERTED_DIR, os.path.basename(docx_path).replace(".docx", ".md"))
        doc = Document(docx_path)
        md_content = []
        
        list_level = 0
        in_table = False
        table_data = []
        
        def get_heading_level(style_name):
            """Extract heading level from style name"""
            if style_name.startswith('Heading'):
                try:
                    return int(style_name[-1])
                except ValueError:
                    return 1
            return 0
        
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if not text and not in_table:  # Skip empty paragraphs unless in table
                md_content.append("")
                continue
            
            # Handle headings
            heading_level = get_heading_level(paragraph.style.name)
            if heading_level > 0:
                md_content.append(f"\n{'#' * heading_level} {text}\n")
                continue
            
            # Handle lists
            if paragraph.style.name.startswith('List'):
                current_level = int(paragraph.style.name[-1]) if paragraph.style.name[-1].isdigit() else 1
                prefix = '  ' * (current_level - 1) + '* '
                md_content.append(f"{prefix}{text}")
                continue
            
            # Process paragraph content with formatting
            formatted_text = ''
            for run in paragraph.runs:
                content = run.text
                if content.strip():
                    # Apply formatting
                    if run.bold and run.italic:
                        content = f"***{content}***"
                    elif run.bold:
                        content = f"**{content}**"
                    elif run.italic:
                        content = f"*{content}*"
                    if run.underline:
                        content = f"_{content}_"
                    if run.strike:
                        content = f"~~{content}~~"
                    
                    formatted_text += content
            
            if formatted_text:
                md_content.append(formatted_text)
        
        # Handle tables
        for table in doc.tables:
            md_content.append("\n")  # Add spacing before table
            
            # Get max column widths
            col_widths = [0] * len(table.columns)
            for row in table.rows:
                for i, cell in enumerate(row.cells):
                    col_widths[i] = max(col_widths[i], len(cell.text.strip()))
            
            # Create header row
            header = "|"
            for i, cell in enumerate(table.rows[0].cells):
                header += f" {cell.text.strip():{col_widths[i]}} |"
            md_content.append(header)
            
            # Create separator
            separator = "|"
            for width in col_widths:
                separator += f" {'-' * width} |"
            md_content.append(separator)
            
            # Create data rows
            for row in table.rows[1:]:
                data_row = "|"
                for i, cell in enumerate(row.cells):
                    data_row += f" {cell.text.strip():{col_widths[i]}} |"
                md_content.append(data_row)
            
            md_content.append("\n")  # Add spacing after table
        
        # Write to file
        with open(md_path, "w", encoding="utf-8") as f:
            f.write("\n".join(md_content))
        
        return md_path
        
    except Exception as e:
        print(f"Error in DOCX to Markdown conversion: {str(e)}")
        traceback_info = traceback.format_exc()
        print(f"Traceback: {traceback_info}")
        raise Exception(f"DOCX to Markdown conversion failed: {str(e)}")

# ✅ DOCX to Markdown
def convert_docx_to_markdown(docx_path):
    md_path = os.path.join(CONVERTED_DIR, os.path.basename(docx_path).replace(".docx", ".md"))
    doc = Document(docx_path)
    md_content = "\n\n".join([p.text for p in doc.paragraphs])
    with open(md_path, "w", encoding="utf-8") as f:
        f.write(md_content)
    return md_path

# ✅ Excel to CSV
def convert_xlsx_to_csv(xlsx_path):
    csv_path = os.path.join(CONVERTED_DIR, os.path.basename(xlsx_path).replace(".xlsx", ".csv"))
    df = pd.read_excel(xlsx_path)
    df.to_csv(csv_path, index=False)
    return csv_path

# ✅ Excel to JSON
def convert_xlsx_to_json(xlsx_path):
    json_path = os.path.join(CONVERTED_DIR, os.path.basename(xlsx_path).replace(".xlsx", ".json"))
    df = pd.read_excel(xlsx_path)
    df.to_json(json_path, orient="records", indent=4)
    return json_path

# ✅ Excel to XML
def convert_xlsx_to_xml(xlsx_path):
    try:
        print(f"Starting Excel to XML conversion: {xlsx_path}")
        xml_path = os.path.join(CONVERTED_DIR, os.path.basename(xlsx_path).replace(".xlsx", ".xml"))
        
        # Read Excel file
        df = pd.read_excel(xlsx_path)
        
        # Convert to XML with proper root element
        records = df.to_dict(orient="records")
        
        # Create XML document with a single root
        xml_data = '<?xml version="1.0" encoding="UTF-8"?>\n'
        xml_data += '<data>\n'
        
        # Add each record as a row element
        for i, record in enumerate(records):
            xml_data += f'  <row id="{i+1}">\n'
            for key, value in record.items():
                # Handle different data types
                if pd.isna(value):
                    value = ""
                elif isinstance(value, (int, float)):
                    value = str(value)
                else:
                    # Escape XML special characters
                    value = str(value).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                
                xml_data += f'    <{key}>{value}</{key}>\n'
            xml_data += '  </row>\n'
        
        xml_data += '</data>'
        
        # Write XML file
        with open(xml_path, 'w', encoding='utf-8') as f:
            f.write(xml_data)
            
        print(f"Excel to XML conversion completed successfully: {xml_path}")
        return xml_path
        
    except Exception as e:
        print(f"Error in Excel to XML conversion: {str(e)}")
        traceback_info = traceback.format_exc()
        print(f"Traceback: {traceback_info}")
        raise Exception(f"Excel to XML conversion failed: {str(e)}")

# ✅ PowerPoint to PDF
async def convert_pptx_to_pdf(pptx_path, task_id=None):
    try:
        if task_id:
            conversion_progress[task_id] = {"progress": 0, "status": "Starting PowerPoint to PDF conversion"}
        
        print(f"Starting PowerPoint to PDF conversion: {pptx_path}")
        pdf_path = os.path.join(CONVERTED_DIR, os.path.basename(pptx_path).replace(".pptx", ".pdf"))
        print(f"Output will be saved to: {pdf_path}")

        # Skip desktop-only methods when running on a server
        if not running_on_server:
            # Try using win32com on Windows first (best formatting)
            try:
                import platform
                if platform.system() == 'Windows':
                    import win32com.client
                    if task_id:
                        conversion_progress[task_id] = {"progress": 10, "status": "Using PowerPoint native conversion"}

                    # Use absolute paths to avoid issues
                    abs_pptx_path = os.path.abspath(pptx_path)
                    abs_pdf_path = os.path.abspath(pdf_path)
                    
                    # Initialize PowerPoint with error handling
                    try:
                        # Try to kill any existing PowerPoint processes first to avoid conflicts
                        try:
                            import subprocess
                            subprocess.call("taskkill /f /im POWERPNT.EXE", shell=True)
                        except:
                            pass
                        
                        # Initialize PowerPoint with timeout
                        import threading
                        import time
                        
                        powerpoint = None
                        init_error = None
                        
                        def init_powerpoint():
                            nonlocal powerpoint, init_error
                            try:
                                powerpoint = win32com.client.Dispatch("Powerpoint.Application")
                                # Try to make PowerPoint invisible (two different methods)
                                try:
                                    # First method: set Visible property to 0
                                    powerpoint.Visible = 0
                                except Exception as vis_error:
                                    print(f"Could not set PowerPoint visibility (method 1): {str(vis_error)}")
                                    try:
                                        # Second method: try the WindowState property
                                        powerpoint.WindowState = 2  # Minimized
                                    except:
                                        pass
                                
                                try:
                                    # Additional approach: Use Windows API to hide window if possible
                                    import ctypes
                                    user32 = ctypes.windll.user32
                                    # Try to find PowerPoint window
                                    hwnd = user32.FindWindowW(None, "PowerPoint")
                                    if hwnd:
                                        # Hide the window: SW_HIDE = 0
                                        user32.ShowWindow(hwnd, 0)
                                except:
                                    pass
                            except Exception as e:
                                init_error = e
                        
                        # Start PowerPoint initialization in a separate thread with timeout
                        init_thread = threading.Thread(target=init_powerpoint)
                        init_thread.daemon = True
                        init_thread.start()
                        init_thread.join(timeout=30)  # Wait up to 30 seconds
                        
                        if init_error:
                            raise init_error
                        
                        if not powerpoint:
                            raise Exception("PowerPoint initialization timed out")
                        
                        # Set higher quality export options if possible
                        try:
                            powerpoint.ActivePrinter = "Microsoft Print to PDF"
                        except:
                            print("Could not set PDF printer, using default")
                        
                        # Open the presentation with proper error handling
                        try:
                            # Try to open with timeout
                            deck = None
                            open_error = None
                            
                            def open_presentation():
                                nonlocal deck, open_error
                                try:
                                    deck = powerpoint.Presentations.Open(abs_pptx_path, WithWindow=False, ReadOnly=True)
                                except Exception as e:
                                    open_error = e
                            
                            # Start presentation opening in a separate thread with timeout
                            open_thread = threading.Thread(target=open_presentation)
                            open_thread.daemon = True
                            open_thread.start()
                            open_thread.join(timeout=60)  # Wait up to 60 seconds
                            
                            if open_error:
                                raise open_error
                            
                            if not deck:
                                raise Exception("Opening presentation timed out")
                            
                            # Try different export methods in order of quality
                            export_success = False
                            
                            # Method 1: ExportAsFixedFormat with high quality settings
                            try:
                                print("Trying ExportAsFixedFormat with high quality settings")
                                deck.ExportAsFixedFormat(
                                    abs_pdf_path, 
                                    FixedFormatType=2,  # PDF
                                    Intent=2,  # Print quality
                                    FrameSlides=False,
                                    OptimizeForPrinter=True,
                                    UseISO19005_1=False,
                                    IncludeDocProperties=True,
                                    KeepIRMSettings=True,
                                    DocStructureTags=True,
                                    BitmapMissingFonts=True,
                                    UseNewProcess=False
                                )
                                export_success = True
                            except Exception as e:
                                print(f"Method 1 failed: {str(e)}")
                            
                            # Method 2: PrintOut to Microsoft Print to PDF
                            if not export_success:
                                try:
                                    print("Trying PrintOut to PDF")
                                    # Try to set the printer to Microsoft Print to PDF
                                    try:
                                        powerpoint.ActivePrinter = "Microsoft Print to PDF"
                                    except:
                                        print("Could not set PDF printer")
                                    
                                    # Print to PDF
                                    deck.PrintOut(
                                        PrintToFile=True,
                                        OutputFileName=abs_pdf_path
                                    )
                                    export_success = True
                                except Exception as e:
                                    print(f"Method 2 failed: {str(e)}")
                            
                            # Method 3: Basic SaveAs
                            if not export_success:
                                try:
                                    print("Trying basic SaveAs")
                                    deck.SaveAs(abs_pdf_path, 32)  # 32 is the PDF format code
                                    export_success = True
                                except Exception as e:
                                    print(f"Method 3 failed: {str(e)}")
                            
                            # Close PowerPoint properly
                            try:
                                deck.Close()
                                powerpoint.Quit()
                            except:
                                print("Error closing PowerPoint, trying to force quit")
                                try:
                                    import subprocess
                                    subprocess.call("taskkill /f /im POWERPNT.EXE", shell=True)
                                except:
                                    pass
                            
                            if export_success and os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 0:
                                print("Successfully converted using PowerPoint native conversion")
                                if task_id:
                                    conversion_progress[task_id] = {"progress": 100, "status": "Conversion complete"}
                                return pdf_path
                            else:
                                raise Exception("All PowerPoint export methods failed")
                        except Exception as e:
                            print(f"Error during presentation processing: {str(e)}")
                            # Clean up
                            if 'deck' in locals() and deck:
                                try:
                                    deck.Close()
                                except:
                                    pass
                            if powerpoint:
                                try:
                                    powerpoint.Quit()
                                except:
                                    pass
                                try:
                                    import subprocess
                                    subprocess.call("taskkill /f /im POWERPNT.EXE", shell=True)
                                except:
                                    pass
                            raise
                    except Exception as e:
                        print(f"Error initializing PowerPoint: {str(e)}")
                        if 'powerpoint' in locals() and powerpoint:
                            try:
                                powerpoint.Quit()
                            except:
                                pass
                            try:
                                import subprocess
                                subprocess.call("taskkill /f /im POWERPNT.EXE", shell=True)
                            except:
                                pass
                        raise
            except Exception as e:
                print(f"Native PowerPoint conversion failed: {str(e)}")
                print("Falling back to alternative method...")

            # Try using comtypes on Windows as fallback
            try:
                if platform.system() == 'Windows':
                    # Check if comtypes is available
                    comtypes_available = False
                    try:
                        import comtypes.client
                        comtypes_available = True
                    except ImportError:
                        print("comtypes module not found. Install with 'pip install comtypes' for better PowerPoint conversion on Windows.")
                        
                    if comtypes_available:
                        if task_id:
                            conversion_progress[task_id] = {"progress": 10, "status": "Using comtypes conversion"}

                        # Use absolute paths
                        abs_pptx_path = os.path.abspath(pptx_path)
                        abs_pdf_path = os.path.abspath(pdf_path)
                        
                        # Initialize PowerPoint with error handling
                        try:
                            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
                            # Don't try to hide PowerPoint as it can cause errors in some environments
                            # powerpoint.Visible = False
                            
                            # Open the presentation with proper error handling
                            try:
                                deck = powerpoint.Presentations.Open(abs_pptx_path, WithWindow=False)
                                
                                # Try to use enhanced PDF export if available
                                try:
                                    # Constants for PDF export
                                    ppFixedFormatTypePDF = 2
                                    ppFixedFormatIntentPrint = 2
                                    
                                    deck.ExportAsFixedFormat(
                                        abs_pdf_path,
                                        ppFixedFormatTypePDF,
                                        Intent=ppFixedFormatIntentPrint,
                                        OptimizeForPrinter=True,
                                        IncludeDocProperties=True,
                                        DocStructureTags=True,
                                        BitmapMissingFonts=True
                                    )
                                except:
                                    # Fallback to basic save
                                    print("Advanced PDF export failed, using basic SaveAs")
                                    deck.SaveAs(abs_pdf_path, 32)  # 32 is the PDF format code
                                
                                deck.Close()
                                powerpoint.Quit()
                                
                                if os.path.exists(pdf_path) and os.path.getsize(pdf_path) > 0:
                                    print("Successfully converted using comtypes")
                                    if task_id:
                                        conversion_progress[task_id] = {"progress": 100, "status": "Conversion complete"}
                                    return pdf_path
                            except Exception as e:
                                print(f"Error opening presentation with comtypes: {str(e)}")
                                if 'deck' in locals() and deck:
                                    try:
                                        deck.Close()
                                    except:
                                        pass
                                if 'powerpoint' in locals() and powerpoint:
                                    try:
                                        powerpoint.Quit()
                                    except:
                                        pass
                                raise
                        except Exception as e:
                            print(f"Error initializing PowerPoint with comtypes: {str(e)}")
                            if 'powerpoint' in locals() and powerpoint:
                                try:
                                    powerpoint.Quit()
                                except:
                                    pass
                            raise
            except Exception as e:
                print(f"Comtypes conversion failed: {str(e)}")
                print("Falling back to python-pptx conversion...")
        else:
            print("Skipping Windows-only conversion methods in server environment")

        # Try using LibreOffice if available (Linux/Mac/Windows)
        try:
            if task_id:
                conversion_progress[task_id] = {"progress": 10, "status": "Trying LibreOffice conversion"}
            
            # Check for LibreOffice in common locations
            libreoffice_paths = [
                # Windows
                r"C:\Program Files\LibreOffice\program\soffice.exe",
                r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                # Linux
                "/usr/bin/libreoffice",
                "/usr/bin/soffice",
                # macOS
                "/Applications/LibreOffice.app/Contents/MacOS/soffice"
            ]
            
            libreoffice_path = None
            for path in libreoffice_paths:
                if os.path.exists(path):
                    libreoffice_path = path
                    break
            
            if libreoffice_path:
                print(f"Found LibreOffice at {libreoffice_path}")
                
                # Use absolute paths
                abs_pptx_path = os.path.abspath(pptx_path)
                abs_pdf_path = os.path.abspath(pdf_path)
                
                # Create a temporary directory for conversion to avoid permission issues
                temp_conversion_dir = os.path.join(CONVERTED_DIR, "temp_libreoffice_conversion")
                os.makedirs(temp_conversion_dir, exist_ok=True)
                
                # Copy the PPTX to the temp directory to avoid permission issues
                import shutil
                temp_pptx_path = os.path.join(temp_conversion_dir, os.path.basename(pptx_path))
                shutil.copy2(abs_pptx_path, temp_pptx_path)
                
                # Run LibreOffice headless to convert with enhanced options
                import subprocess
                
                # First try with enhanced PDF export options
                enhanced_cmd = [
                    libreoffice_path,
                    "--headless",
                    "--convert-to", "pdf:writer_pdf_Export:{'ExportBookmarks':1,'Quality':100,'ReduceImageResolution':0,'EmbedStandardFonts':1}",
                    "--outdir", temp_conversion_dir,
                    temp_pptx_path
                ]
                
                try:
                    process = subprocess.Popen(enhanced_cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
                    stdout, stderr = process.communicate(timeout=60)  # Add timeout to prevent hanging
                    
                    # Check if conversion was successful
                    temp_pdf_path = os.path.join(temp_conversion_dir, os.path.basename(temp_pptx_path).replace(".pptx", ".pdf"))
                    if process.returncode == 0 and os.path.exists(temp_pdf_path) and os.path.getsize(temp_pdf_path) > 0:
                        # Copy the PDF to the final destination
                        shutil.copy2(temp_pdf_path, abs_pdf_path)
                        print("Successfully converted using LibreOffice with enhanced options")
                        
                        # Clean up temp directory
                        try:
                            shutil.rmtree(temp_conversion_dir)
                        except:
                            print("Warning: Failed to clean up temp LibreOffice directory")
                        
                        if task_id:
                            conversion_progress[task_id] = {"progress": 100, "status": "Conversion complete"}
                        return pdf_path
                except subprocess.TimeoutExpired:
                    print("LibreOffice conversion timed out with enhanced options")
                    # Kill the process if it's still running
                    try:
                        process.kill()
                    except:
                        pass
                
                # If enhanced conversion failed, try with basic options
                basic_cmd = [
                    libreoffice_path,
                    "--headless",
                    "--convert-to", "pdf",
                    "--outdir", temp_conversion_dir,
                    temp_pptx_path
                ]
                
                try:
                    process = subprocess.Popen(basic_cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
                    stdout, stderr = process.communicate(timeout=60)
                    
                    # Check if conversion was successful
                    temp_pdf_path = os.path.join(temp_conversion_dir, os.path.basename(temp_pptx_path).replace(".pptx", ".pdf"))
                    if process.returncode == 0 and os.path.exists(temp_pdf_path) and os.path.getsize(temp_pdf_path) > 0:
                        # Copy the PDF to the final destination
                        shutil.copy2(temp_pdf_path, abs_pdf_path)
                        print("Successfully converted using LibreOffice with basic options")
                        
                        # Clean up temp directory
                        try:
                            shutil.rmtree(temp_conversion_dir)
                        except:
                            print("Warning: Failed to clean up temp LibreOffice directory")
                        
                        if task_id:
                            conversion_progress[task_id] = {"progress": 100, "status": "Conversion complete"}
                        return pdf_path
                except subprocess.TimeoutExpired:
                    print("LibreOffice conversion timed out with basic options")
                    # Kill the process if it's still running
                    try:
                        process.kill()
                    except:
                        pass
                
                # Clean up temp directory if conversion failed
                try:
                    shutil.rmtree(temp_conversion_dir)
                except:
                    print("Warning: Failed to clean up temp LibreOffice directory")
                
                print("LibreOffice conversion failed, trying alternative methods")
            else:
                print("LibreOffice not found, using Python fallback")
        except Exception as e:
            print(f"LibreOffice conversion failed: {str(e)}")
            print("Falling back to python-pptx with enhanced image quality...")

        # Fallback to python-pptx with enhanced image quality
        if task_id:
            conversion_progress[task_id] = {"progress": 10, "status": "Using enhanced image-based conversion"}

        print("Using enhanced image-based conversion")
        prs = Presentation(pptx_path)
        
        # Create temporary directory for slide images
        temp_dir = os.path.join(CONVERTED_DIR, "temp_slides")
        os.makedirs(temp_dir, exist_ok=True)
        
        total_slides = len(prs.slides)
        images = []

        # Get presentation dimensions
        slide_width = int(prs.slide_width * 96 / 914400)  # Convert EMU to pixels at 96 DPI
        slide_height = int(prs.slide_height * 96 / 914400)
        
        # Use higher resolution for better quality
        scale_factor = 3  # Increase resolution by 3x for better quality
        img_width = slide_width * scale_factor
        img_height = slide_height * scale_factor

        # Try to detect background color from first slide
        background_color = (255, 255, 255)  # Default white
        try:
            if prs.slides and hasattr(prs.slides[0], 'background'):
                bg = prs.slides[0].background
                if hasattr(bg, 'fill') and hasattr(bg.fill, 'fore_color'):
                    if hasattr(bg.fill.fore_color, 'rgb'):
                        r = bg.fill.fore_color.rgb[0]
                        g = bg.fill.fore_color.rgb[1]
                        b = bg.fill.fore_color.rgb[2]
                        background_color = (r, g, b)
        except:
            pass  # Use default white if detection fails
            
        # Try alternative approach: export to individual PPTX files and convert
        try:
            # Check if we're on Windows and can use PowerPoint
            import platform
            if platform.system() == 'Windows':
                try:
                    import win32com.client
                    
                    # Create a directory for individual slides
                    slides_dir = os.path.join(temp_dir, "individual_slides")
                    os.makedirs(slides_dir, exist_ok=True)
                    
                    # Ensure the directory has proper permissions
                    try:
                        # Make sure the directory is writeable
                        os.chmod(slides_dir, stat.S_IRWXU | stat.S_IRWXG | stat.S_IRWXO)  # Full permissions
                    except Exception as e:
                        print(f"Warning: Could not set permissions on directory: {str(e)}")
                    
                    # Export each slide as a separate PPTX
                    for i, slide in enumerate(prs.slides):
                        if task_id:
                            progress = 10 + (i / total_slides * 30)
                            conversion_progress[task_id] = {
                                "progress": progress,
                                "status": f"Preparing slide {i+1}/{total_slides}"
                            }
                        
                        # Create a new presentation with just this slide
                        temp_prs = Presentation()
                        # Copy slide size and properties
                        temp_prs.slide_width = prs.slide_width
                        temp_prs.slide_height = prs.slide_height
                        
                        # Add a blank slide
                        blank_slide_layout = temp_prs.slide_layouts[6]  # Blank layout
                        temp_slide = temp_prs.slides.add_slide(blank_slide_layout)
                        
                        # Save the temporary presentation
                        temp_pptx_path = os.path.join(slides_dir, f"slide_{i+1:03d}.pptx")
                        temp_prs.save(temp_pptx_path)
                    
                    # Initialize PowerPoint once for all slides
                    powerpoint = win32com.client.Dispatch("Powerpoint.Application")
                    
                    # Try to make PowerPoint invisible (two different methods)
                    try:
                        # First method: set Visible property to 0
                        powerpoint.Visible = 0
                    except Exception as vis_error:
                        print(f"Could not set PowerPoint visibility (method 1): {str(vis_error)}")
                        try:
                            # Second method: try the WindowState property
                            powerpoint.WindowState = 2  # Minimized
                        except:
                            pass
                    
                    try:
                        # Additional approach: Use Windows API to hide window if possible
                        import ctypes
                        user32 = ctypes.windll.user32
                        # Try to find PowerPoint window
                        hwnd = user32.FindWindowW(None, "PowerPoint")
                        if hwnd:
                            # Hide the window: SW_HIDE = 0
                            user32.ShowWindow(hwnd, 0)
                    except:
                        pass
                    
                    # Open the main presentation
                    main_deck = powerpoint.Presentations.Open(os.path.abspath(pptx_path), WithWindow=False, ReadOnly=True)
                    
                    # Process each slide
                    for i in range(total_slides):
                        if task_id:
                            progress = 40 + (i / total_slides * 40)
                            conversion_progress[task_id] = {
                                "progress": progress,
                                "status": f"Converting slide {i+1}/{total_slides} to image"
                            }
                        
                        # Get the individual slide PPTX
                        temp_pptx_path = os.path.join(slides_dir, f"slide_{i+1:03d}.pptx")
                        temp_pdf_path = os.path.join(slides_dir, f"slide_{i+1:03d}.pdf")
                        
                        # Open the temporary presentation
                        temp_deck = powerpoint.Presentations.Open(os.path.abspath(temp_pptx_path), WithWindow=False)
                        
                        # Copy the slide from main presentation to temp
                        main_deck.Slides(i+1).Copy()
                        temp_deck.Slides(1).Delete()  # Delete the blank slide
                        temp_deck.Paste()
                        
                        # Export as PDF
                        try:
                            temp_deck.ExportAsFixedFormat(
                                os.path.abspath(temp_pdf_path),
                                FixedFormatType=2,  # PDF
                                Intent=2,  # Print quality
                                OptimizeForPrinter=True
                            )
                        except:
                            # Fallback to basic save
                            temp_deck.SaveAs(os.path.abspath(temp_pdf_path), 32)
                        
                        # Close the temp presentation
                        temp_deck.Close()
                        
                        # Convert PDF to high-quality image
                        try:
                            from pdf2image import convert_from_path
                            pdf_images = convert_from_path(temp_pdf_path, dpi=300)
                            if pdf_images:
                                # Save the image
                                img_path = os.path.join(temp_dir, f"slide_{i+1:03d}.png")
                                pdf_images[0].save(img_path, "PNG", quality=100, optimize=True, dpi=(300, 300))
                                images.append(pdf_images[0])
                                print(f"Successfully converted slide {i+1} using PowerPoint and pdf2image")
                        except Exception as pdf_error:
                            print(f"PDF to image conversion failed: {str(pdf_error)}")
                    
                    # Close the main presentation and PowerPoint
                    main_deck.Close()
                    powerpoint.Quit()
                    
                    # If we have all images, we can skip the manual rendering
                    if len(images) == total_slides:
                        print("Successfully converted all slides using PowerPoint export")
                        # Skip to PDF creation
                        if task_id:
                            conversion_progress[task_id] = {"progress": 80, "status": "Creating PDF file"}
                        # Continue to PDF creation below
                    else:
                        print(f"Only converted {len(images)}/{total_slides} slides, falling back to manual rendering")
                        # Clear images to force manual rendering
                        images = []
                except Exception as win32_error:
                    print(f"PowerPoint slide export failed: {str(win32_error)}")
                    # Clear images to force manual rendering
                    images = []
        except Exception as export_error:
            print(f"Alternative slide export failed: {str(export_error)}")
            # Continue with manual rendering

        # If we don't have all images, use manual rendering
        if len(images) != total_slides:
            # Process each slide manually
            for i, slide in enumerate(prs.slides):
                if task_id:
                    progress = 10 + (i / total_slides * 60)
                    conversion_progress[task_id] = {
                        "progress": progress,
                        "status": f"Processing slide {i+1}/{total_slides}"
                    }

                # Create high-resolution slide image with correct aspect ratio
                img = Image.new('RGB', (img_width, img_height), background_color)
                draw = ImageDraw.Draw(img)
                
                # Try to use high-quality fonts
                title_font_size = int(36 * scale_factor)
                body_font_size = int(24 * scale_factor)
                small_font_size = int(18 * scale_factor)
                
                # Find system fonts
                try:
                    # Try to find system fonts
                    system_fonts = [
                        "arial.ttf", "Arial.ttf",
                        "calibri.ttf", "Calibri.ttf",
                        "segoeui.ttf", "SegoeUI.ttf",
                        "times.ttf", "Times New Roman.ttf",
                        "DejaVuSans.ttf", "LiberationSans-Regular.ttf"
                    ]
                    
                    # Common font directories
                    font_dirs = [
                        # Windows
                        r"C:\Windows\Fonts",
                        # macOS
                        "/Library/Fonts",
                        "/System/Library/Fonts",
                        # Linux
                        "/usr/share/fonts",
                        "/usr/local/share/fonts"
                    ]
                    
                    # Find a usable font
                    title_font = None
                    body_font = None
                    small_font = None
                    
                    for font_dir in font_dirs:
                        if os.path.exists(font_dir):
                            for font_name in system_fonts:
                                font_path = os.path.join(font_dir, font_name)
                                if os.path.exists(font_path):
                                    title_font = ImageFont.truetype(font_path, title_font_size)
                                    body_font = ImageFont.truetype(font_path, body_font_size)
                                    small_font = ImageFont.truetype(font_path, small_font_size)
                                    break
                            if title_font:
                                break
                    
                    if not title_font:
                        raise Exception("No system fonts found")
                        
                except:
                    # Fallback to default font
                    title_font = ImageFont.load_default()
                    body_font = title_font
                    small_font = title_font
                    print("Using default font as fallback")

                # Process slide content with better positioning
                # First, try to extract slide background
                try:
                    if hasattr(slide, 'background') and hasattr(slide.background, 'fill'):
                        bg_fill = slide.background.fill
                        if hasattr(bg_fill, 'fore_color') and hasattr(bg_fill.fore_color, 'rgb'):
                            r, g, b = bg_fill.fore_color.rgb
                            img = Image.new('RGB', (img_width, img_height), (r, g, b))
                            draw = ImageDraw.Draw(img)
                except:
                    pass  # Keep default background if extraction fails
                
                # Process shapes in z-order (back to front)
                for shape in slide.shapes:
                    try:
                        # Convert shape position from EMU to pixels
                        left = int(shape.left * img_width / prs.slide_width)
                        top = int(shape.top * img_height / prs.slide_height)
                        width = int(shape.width * img_width / prs.slide_width)
                        height = int(shape.height * img_height / prs.slide_height)
                        
                        # Handle different shape types
                        if shape.shape_type == 1:  # Auto Shape
                            # Draw rectangle as placeholder
                            draw.rectangle([left, top, left + width, top + height], 
                                          outline=(200, 200, 200), width=2)
                        
                        # Handle text
                        if hasattr(shape, 'text') and shape.text.strip():
                            # Determine text style
                            is_title = any(word in shape.name.lower() for word in ['title', 'header'])
                            font = title_font if is_title else body_font
                            
                            # Get text color if available
                            text_color = (0, 0, 0)  # Default black
                            try:
                                if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                                    p = shape.text_frame.paragraphs[0]
                                    if p.runs and hasattr(p.runs[0].font, 'color') and p.runs[0].font.color.rgb:
                                        r, g, b = p.runs[0].font.color.rgb
                                        text_color = (r, g, b)
                            except:
                                pass
                            
                            # Word wrap text to fit shape width
                            text = shape.text.strip()
                            words = text.split()
                            lines = []
                            current_line = []
                            
                            for word in words:
                                test_line = ' '.join(current_line + [word])
                                text_width = draw.textlength(test_line, font=font)
                                
                                if text_width < width - 10:  # Leave small margin
                                    current_line.append(word)
                                else:
                                    if current_line:
                                        lines.append(' '.join(current_line))
                                    current_line = [word]
                            
                            if current_line:
                                lines.append(' '.join(current_line))
                            
                            # Draw text with proper positioning
                            y_offset = top
                            for line in lines:
                                draw.text((left + 5, y_offset), line, fill=text_color, font=font)
                                y_offset += int(font.size * 1.2)  # Line spacing
                        
                        # Try to handle images
                        if hasattr(shape, 'image') and shape.shape_type == 13:  # Picture
                            try:
                                # Extract image data if available
                                image_data = None
                                
                                # Method 1: Try to extract from blip element
                                if hasattr(shape, '_element') and hasattr(shape._element, 'blip'):
                                    try:
                                        image_stream = io.BytesIO(shape._element.blip.embed.blob)
                                        shape_img = Image.open(image_stream)
                                        image_data = shape_img
                                    except Exception as e:
                                        print(f"Method 1 image extraction failed: {str(e)}")
                                
                                # Method 2: Try alternative extraction via part relationship
                                if not image_data and hasattr(shape, 'part'):
                                    try:
                                        # Get relationship ID
                                        rId = shape._element.blip_rId
                                        image_part = shape.part.related_parts[rId]
                                        image_stream = io.BytesIO(image_part.blob)
                                        shape_img = Image.open(image_stream)
                                        image_data = shape_img
                                    except Exception as e:
                                        print(f"Method 2 image extraction failed: {str(e)}")
                                
                                # Method 3: Try direct XML extraction
                                if not image_data:
                                    try:
                                        from lxml import etree
                                        
                                        # Get the relationship ID from the shape element
                                        xfrm = shape._element.xpath('.//p:blipFill/a:blip', 
                                                                  namespaces={
                                                                      'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
                                                                      'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
                                                                  })
                                        
                                        if xfrm and 'embed' in xfrm[0].attrib:
                                            rId = xfrm[0].attrib['embed']
                                            image_part = shape.part.related_parts[rId]
                                            image_stream = io.BytesIO(image_part.blob)
                                            shape_img = Image.open(image_stream)
                                            image_data = shape_img
                                    except Exception as e:
                                        print(f"Method 3 image extraction failed: {str(e)}")
                                
                                # If we have image data, process and paste it
                                if image_data:
                                    # Apply high-quality resizing
                                    try:
                                        # Use LANCZOS for high-quality downsampling
                                        shape_img = image_data.resize((width, height), Image.LANCZOS)
                                        
                                        # Create a mask for transparency if needed
                                        if shape_img.mode == 'RGBA':
                                            # Create a white background
                                            bg = Image.new('RGBA', shape_img.size, (255, 255, 255, 255))
                                            # Paste the image using itself as mask
                                            bg.paste(shape_img, (0, 0), shape_img)
                                            shape_img = bg.convert('RGB')
                                        elif shape_img.mode != 'RGB':
                                            shape_img = shape_img.convert('RGB')
                                        
                                        # Paste image onto slide
                                        img.paste(shape_img, (left, top))
                                        print(f"Successfully pasted image in slide {i+1}")
                                    except Exception as resize_error:
                                        print(f"Error resizing image: {str(resize_error)}")
                                        # Draw placeholder with error message
                                        draw.rectangle([left, top, left + width, top + height], 
                                                      outline=(200, 0, 0), width=2)
                                        draw.text((left + 10, top + 10), "Image Error", fill=(200, 0, 0), font=small_font)
                                else:
                                    # Draw placeholder if image extraction fails
                                    draw.rectangle([left, top, left + width, top + height], 
                                                  outline=(150, 150, 150), width=2)
                                    draw.line([left, top, left + width, top + height], fill=(150, 150, 150), width=1)
                                    draw.line([left + width, top, left, top + height], fill=(150, 150, 150), width=1)
                                    draw.text((left + 10, top + 10), "Image", fill=(150, 150, 150), font=small_font)
                            except Exception as img_error:
                                print(f"Image handling error: {str(img_error)}")
                                # Draw placeholder if image extraction fails
                                draw.rectangle([left, top, left + width, top + height], 
                                              outline=(150, 150, 150), width=2)
                                draw.text((left + 10, top + 10), "Image Error", fill=(150, 150, 150), font=small_font)
                    
                    except Exception as shape_error:
                        print(f"Error processing shape: {str(shape_error)}")
                        continue  # Skip problematic shapes

                # Add slide number with subtle styling
                slide_number_text = f"{i+1}/{total_slides}"
                text_width = draw.textlength(slide_number_text, font=small_font)
                draw.text(
                    (img_width - text_width - 20, img_height - small_font.size - 10),
                    slide_number_text,
                    fill=(128, 128, 128),
                    font=small_font
                )

                # Save high-quality image
                img_path = os.path.join(temp_dir, f"slide_{i+1:03d}.png")
                img.save(img_path, "PNG", quality=100, optimize=True, dpi=(300, 300))
                images.append(img)

        # Create PDF with high quality
        if task_id:
            conversion_progress[task_id] = {"progress": 80, "status": "Creating PDF file"}

        print(f"Creating PDF with {len(images)} slides")
        if images:
            # Try using img2pdf for highest quality PDF creation
            try:
                import img2pdf
                
                # Save images to temporary files with consistent naming for sorting
                temp_image_files = []
                for i, img in enumerate(images):
                    temp_img_path = os.path.join(temp_dir, f"slide_{i+1:03d}.png")
                    img.save(temp_img_path, "PNG", quality=100, optimize=True, dpi=(300, 300))
                    temp_image_files.append(temp_img_path)
                
                # Sort files to ensure correct order
                temp_image_files.sort()
                
                # Create PDF with img2pdf for highest quality
                with open(pdf_path, "wb") as f:
                    # Use simpler layout approach - img2pdf.PageSize might not be available in all versions
                    f.write(img2pdf.convert(temp_image_files))
                
                print("Created PDF using img2pdf (highest quality)")
                
            except Exception as img2pdf_error:
                print(f"img2pdf PDF creation failed: {str(img2pdf_error)}")
                print("Falling back to reportlab PDF creation")
                
                # Try using reportlab for better PDF quality if available
                try:
                    from reportlab.lib.pagesizes import landscape
                    from reportlab.pdfgen import canvas
                    from reportlab.lib.utils import ImageReader
                    
                    # Calculate page size based on first image
                    page_width, page_height = images[0].size
                    
                    # Create PDF with reportlab
                    c = canvas.Canvas(pdf_path, pagesize=(page_width, page_height))
                    
                    for i, img in enumerate(images):
                        if i > 0:  # First page is already created
                            c.showPage()
                        
                        # Save image to temporary file for reportlab
                        temp_img_path = os.path.join(temp_dir, f"temp_slide_{i+1:03d}.png")
                        img.save(temp_img_path, "PNG", quality=100)
                        
                        # Add image to PDF with precise positioning
                        c.drawImage(
                            ImageReader(temp_img_path),
                            0, 0,
                            width=page_width,
                            height=page_height,
                            mask=None,
                            preserveAspectRatio=True,
                            anchor='c'
                        )
                    
                    c.save()
                    print("Created PDF using reportlab")
                    
                except Exception as reportlab_error:
                    print(f"Reportlab PDF creation failed: {str(reportlab_error)}")
                    print("Falling back to PIL PDF creation")
                    
                    # Fallback to PIL PDF creation
                    try:
                        # Save images with higher quality first
                        high_quality_images = []
                        for img in images:
                            # Ensure image is in RGB mode
                            if img.mode != 'RGB':
                                img = img.convert('RGB')
                            high_quality_images.append(img)
                        
                        # Create PDF with PIL
                        high_quality_images[0].save(
                            pdf_path,
                            save_all=True,
                            append_images=high_quality_images[1:],
                            resolution=300.0,
                            quality=100
                        )
                        print("Created PDF using PIL")
                    except Exception as pil_error:
                        print(f"PIL PDF creation failed: {str(pil_error)}")
                        raise Exception("All PDF creation methods failed")

        # Clean up
        if task_id:
            conversion_progress[task_id] = {"progress": 90, "status": "Cleaning up"}

        try:
            # Use a more robust cleanup approach
            import shutil
            import time
            
            # Give Windows time to release file handles
            time.sleep(1)
            
            # Try to remove files first
            for file in os.listdir(temp_dir):
                file_path = os.path.join(temp_dir, file)
                try:
                    if os.path.isfile(file_path):
                        os.remove(file_path)
                    elif os.path.isdir(file_path):
                        shutil.rmtree(file_path, ignore_errors=True)
                except Exception as e:
                    print(f"Warning: Could not remove {file_path}: {str(e)}")
            
            # Try to remove the main temp directory
            try:
                os.rmdir(temp_dir)
            except Exception as e:
                print(f"Warning: Could not remove temp directory: {str(e)}")
                # Don't let this stop the conversion
        except Exception as e:
            print(f"Warning: Failed to clean up temp files: {str(e)}")
            # Continue anyway - don't let cleanup issues stop the conversion

        if task_id:
            conversion_progress[task_id] = {"progress": 100, "status": "Conversion complete"}

        return pdf_path

    except Exception as e:
        print(f"Error in PowerPoint to PDF conversion: {str(e)}")
        traceback_info = traceback.format_exc()
        print(f"Traceback: {traceback_info}")
        
        if task_id:
            conversion_progress[task_id] = {"progress": 100, "status": f"Error: {str(e)}"}
        
        raise Exception(f"PowerPoint to PDF conversion failed: {str(e)}")

# ✅ PowerPoint to Images
def convert_pptx_to_images(pptx_path):
    try:
        print(f"Starting PowerPoint to Images conversion: {pptx_path}")
        
        # Create output directory for images
        img_dir = os.path.join(CONVERTED_DIR, os.path.splitext(os.path.basename(pptx_path))[0])
        os.makedirs(img_dir, exist_ok=True)
        print(f"Images will be saved to directory: {img_dir}")
        
        # Open the presentation
        prs = Presentation(pptx_path)
        print(f"Opened PowerPoint with {len(prs.slides)} slides")
        
        img_paths = []
        
        # Use temporary directory for slide rendering
        temp_dir = os.path.join(CONVERTED_DIR, "temp_slides")
        os.makedirs(temp_dir, exist_ok=True)
        
        # Try using Windows PowerPoint for better quality (if available and not in server environment)
        windows_export_success = False
        if not running_on_server:
            try:
                import platform
                if platform.system() == 'Windows':
                    try:
                        import win32com.client
                        
                        # Create a directory for individual slides
                        slides_dir = os.path.join(temp_dir, "individual_slides")
                        os.makedirs(slides_dir, exist_ok=True)
                        
                        # Ensure the directory has proper permissions
                        try:
                            # Make sure the directory is writeable
                            import stat
                            os.chmod(slides_dir, stat.S_IRWXU | stat.S_IRWXG | stat.S_IRWXO)  # Full permissions
                        except Exception as e:
                            print(f"Warning: Could not set permissions on directory: {str(e)}")
                        
                        # Try to kill any existing PowerPoint processes
                        try:
                            import subprocess
                            subprocess.call("taskkill /f /im POWERPNT.EXE", shell=True)
                        except:
                            pass
                        
                        # Initialize PowerPoint
                        powerpoint = win32com.client.Dispatch("Powerpoint.Application")
                        
                        # Try to make PowerPoint invisible (two different methods)
                        try:
                            # First method: set Visible property to 0
                            powerpoint.Visible = 0
                        except Exception as vis_error:
                            print(f"Could not set PowerPoint visibility (method 1): {str(vis_error)}")
                            try:
                                # Second method: try the WindowState property
                                powerpoint.WindowState = 2  # Minimized
                            except:
                                pass
                        
                        try:
                            # Additional approach: Use Windows API to hide window if possible
                            import ctypes
                            user32 = ctypes.windll.user32
                            # Try to find PowerPoint window
                            hwnd = user32.FindWindowW(None, "PowerPoint")
                            if hwnd:
                                # Hide the window: SW_HIDE = 0
                                user32.ShowWindow(hwnd, 0)
                        except:
                            pass
                        
                        # Open the main presentation
                        main_deck = powerpoint.Presentations.Open(os.path.abspath(pptx_path), ReadOnly=True)
                        
                        # Process each slide
                        successful_slides = 0
                        for i in range(len(prs.slides)):
                            try:
                                # Export directly to PNG if possible
                                final_image_path = os.path.join(img_dir, f"slide_{i+1}.png")
                                
                                # Try direct export method
                                try:
                                    # Export the slide directly
                                    main_deck.Slides(i+1).Export(os.path.abspath(final_image_path), "PNG")
                                    img_paths.append(final_image_path)
                                    successful_slides += 1
                                    print(f"Created image for slide {i+1}: {final_image_path}")
                                    continue  # Skip to next slide if successful
                                except Exception as direct_export_error:
                                    print(f"Direct slide export failed for slide {i+1}: {str(direct_export_error)}")
                                
                                # If direct export failed, try creating a single-slide presentation
                                temp_pptx_path = os.path.join(slides_dir, f"slide_{i+1:03d}.pptx")
                                
                                # Create a new presentation with just this slide
                                temp_prs = Presentation()
                                temp_prs.slide_width = prs.slide_width
                                temp_prs.slide_height = prs.slide_height
                                blank_slide_layout = temp_prs.slide_layouts[6]  # Blank layout
                                temp_slide = temp_prs.slides.add_slide(blank_slide_layout)
                                temp_prs.save(temp_pptx_path)
                                
                                # Open the temporary presentation
                                temp_deck = powerpoint.Presentations.Open(os.path.abspath(temp_pptx_path), ReadOnly=False)
                                
                                # Copy the slide from main presentation to temp
                                main_deck.Slides(i+1).Copy()
                                try:
                                    temp_deck.Slides(1).Delete()  # Delete the blank slide
                                except:
                                    pass  # If delete fails, just continue
                                
                                temp_deck.Paste()
                                
                                # Try to save as PNG
                                temp_png_path = os.path.join(slides_dir, f"slide_{i+1:03d}.png")
                                try:
                                    # Try to export directly as PNG
                                    temp_deck.SaveAs(os.path.abspath(temp_png_path), 18)  # ppSaveAsPNG = 18
                                except Exception as png_error:
                                    print(f"PNG export failed for slide {i+1}: {str(png_error)}")
                                    # Try PDF then convert to PNG
                                    try:
                                        temp_pdf_path = os.path.join(slides_dir, f"slide_{i+1:03d}.pdf")
                                        temp_deck.SaveAs(os.path.abspath(temp_pdf_path), 32)  # PDF = 32
                                        
                                        # Convert PDF to PNG
                                        try:
                                            from pdf2image import convert_from_path
                                            pdf_images = convert_from_path(temp_pdf_path, dpi=300)
                                            if pdf_images:
                                                pdf_images[0].save(temp_png_path, "PNG", quality=100, dpi=(300, 300))
                                        except Exception as pdf_error:
                                            print(f"PDF to image conversion failed: {str(pdf_error)}")
                                            raise  # Re-raise to trigger fallback
                                    except:
                                        raise  # Re-raise to trigger fallback
                                
                                # Close the temp presentation
                                try:
                                    temp_deck.Close()
                                except:
                                    pass  # Ignore close errors
                                
                                # Copy the image to the output directory if it exists
                                if os.path.exists(temp_png_path):
                                    import shutil
                                    shutil.copy2(temp_png_path, final_image_path)
                                    img_paths.append(final_image_path)
                                    successful_slides += 1
                                    print(f"Created image for slide {i+1}: {final_image_path}")
                            except Exception as slide_error:
                                print(f"Error processing slide {i+1}: {str(slide_error)}")
                                # Continue with next slide
                        
                        # Close PowerPoint
                        try:
                            main_deck.Close()
                            powerpoint.Quit()
                        except:
                            # Force quit if normal close fails
                            try:
                                import subprocess
                                subprocess.call("taskkill /f /im POWERPNT.EXE", shell=True)
                            except:
                                pass
                        
                        # Check if we got all slides
                        if successful_slides == len(prs.slides):
                            print(f"Successfully created all {successful_slides} slide images")
                            windows_export_success = True
                        else:
                            print(f"Only created {successful_slides} out of {len(prs.slides)} slides")
                            # We'll continue with the fallback method for any missing slides
                    except Exception as win_error:
                        print(f"Windows PowerPoint export failed: {str(win_error)}")
                        # Fall back to the basic method
            except Exception as e:
                print(f"PowerPoint COM automation failed: {str(e)}")
                # Fall back to the basic method
        elif running_on_server:
            print("Skipping Windows PowerPoint export in server environment")
        
        # If Windows export was successful, return the directory
        if windows_export_success and len(img_paths) == len(prs.slides):
            # Clean up temp directory
            try:
                import shutil
                import time
                
                # Give Windows time to release file handles
                time.sleep(1)
                
                # Try to remove the temp directory and its contents
                try:
                    shutil.rmtree(temp_dir, ignore_errors=True)
                except Exception as e:
                    print(f"Warning: Could not remove temp directory: {str(e)}")
            except Exception as e:
                print(f"Warning: Failed to clean up temp files: {str(e)}")
            
            return img_dir
        
        # Fallback: Basic method using python-pptx
        # Only process slides that weren't successfully exported
        existing_slides = set(os.path.basename(path) for path in img_paths)
        
        for i, slide in enumerate(prs.slides):
            # Skip slides that were already successfully exported
            if f"slide_{i+1}.png" in existing_slides:
                continue
                
            # Export slide as image using Python-PPTX
            slide_prs = Presentation()
            # Copy slide dimensions
            slide_prs.slide_width = prs.slide_width
            slide_prs.slide_height = prs.slide_height
            
            slide_layout = slide_prs.slide_layouts[0]  # Use first layout
            slide_copy = slide_prs.slides.add_slide(slide_layout)
            
            # Try to copy more properties
            try:
                # Copy background if possible
                if hasattr(slide, 'background') and hasattr(slide.background, 'fill'):
                    slide_copy.background = slide.background
            except:
                pass
            
            # Copy all shapes from original slide to new slide with better handling
            for shape in slide.shapes:
                try:
                    if hasattr(shape, 'text'):
                        new_shape = slide_copy.shapes.add_textbox(
                            shape.left, shape.top, shape.width, shape.height
                        )
                        new_shape.text = shape.text
                        
                        # Try to copy text properties
                        if hasattr(shape, 'text_frame') and hasattr(new_shape, 'text_frame'):
                            for j, paragraph in enumerate(shape.text_frame.paragraphs):
                                if j < len(new_shape.text_frame.paragraphs):
                                    new_paragraph = new_shape.text_frame.paragraphs[j]
                                    # Copy alignment if available
                                    if hasattr(paragraph, 'alignment') and hasattr(new_paragraph, 'alignment'):
                                        new_paragraph.alignment = paragraph.alignment
                except:
                    # If shape copy fails, continue with the next shape
                    continue
            
            # Save the temporary single-slide presentation
            temp_pptx = os.path.join(temp_dir, f"temp_slide_{i+1}.pptx")
            slide_prs.save(temp_pptx)
            
            # Create a higher quality image for this slide
            img_path = os.path.join(img_dir, f"slide_{i+1}.png")
            
            # Create an image with the correct aspect ratio
            slide_width = int(prs.slide_width * 96 / 914400)  # Convert EMU to pixels
            slide_height = int(prs.slide_height * 96 / 914400)
            
            # Scale up for better quality
            scale_factor = 2
            img_width = slide_width * scale_factor
            img_height = slide_height * scale_factor
            
            # Create a higher quality image
            img = Image.new('RGB', (img_width, img_height), color=(240, 240, 240))
            
            # Add slide number
            try:
                draw = ImageDraw.Draw(img)
                font = None
                
                # Try to use a system font
                try:
                    font_size = int(24 * scale_factor)
                    # Common font directories
                    font_dirs = [
                        # Windows
                        r"C:\Windows\Fonts",
                        # macOS
                        "/Library/Fonts",
                        "/System/Library/Fonts",
                        # Linux
                        "/usr/share/fonts",
                        "/usr/local/share/fonts"
                    ]
                    
                    # Common fonts
                    font_names = [
                        "arial.ttf", "Arial.ttf",
                        "calibri.ttf", "Calibri.ttf",
                        "segoeui.ttf", "SegoeUI.ttf"
                    ]
                    
                    for font_dir in font_dirs:
                        if os.path.exists(font_dir):
                            for font_name in font_names:
                                font_path = os.path.join(font_dir, font_name)
                                if os.path.exists(font_path):
                                    font = ImageFont.truetype(font_path, font_size)
                                    break
                            if font:
                                break
                except:
                    pass
                
                if not font:
                    font = ImageFont.load_default()
                
                # Draw slide number in bottom right
                slide_number_text = f"Slide {i+1}/{len(prs.slides)}"
                text_width = draw.textlength(slide_number_text, font=font)
                draw.text(
                    (img_width - text_width - 20, img_height - font.size - 10),
                    slide_number_text,
                    fill=(100, 100, 100),
                    font=font
                )
            except:
                # If adding slide number fails, continue without it
                pass
            
            # Save the image with high quality
            img.save(img_path, "PNG", quality=100, optimize=True, dpi=(300, 300))
            img_paths.append(img_path)
            
            print(f"Created image for slide {i+1}: {img_path}")
        
        # Clean up temp directory
        try:
            # Use a more robust cleanup approach
            import shutil
            import time
            
            # Give Windows time to release file handles
            time.sleep(1)
            
            # Try to remove the temp directory and its contents
            try:
                shutil.rmtree(temp_dir, ignore_errors=True)
            except Exception as e:
                print(f"Warning: Could not remove temp directory: {str(e)}")
        except Exception as e:
            print(f"Warning: Failed to clean up temp files: {str(e)}")
        
        print(f"Created {len(img_paths)} slide images successfully")
        # Return the directory path instead of the list of image paths
        return img_dir
        
    except Exception as e:
        print(f"Error in PowerPoint to Images conversion: {str(e)}")
        traceback_info = traceback.format_exc()
        print(f"Traceback: {traceback_info}")
        raise Exception(f"PowerPoint to Images conversion failed: {str(e)}")

# ✅ PowerPoint to Video
async def convert_pptx_to_video(pptx_path):
    try:
        print(f"Starting PowerPoint to Video conversion: {pptx_path}")
        
        # First convert PPTX to images
        img_dir = convert_pptx_to_images(pptx_path)
        if not img_dir or not os.path.exists(img_dir):
            raise Exception("Failed to extract images from PowerPoint")
            
        print(f"Successfully extracted images to directory: {img_dir}")
        
        # Get all PNG files from the directory
        img_paths = sorted([
            os.path.join(img_dir, f) for f in os.listdir(img_dir)
            if f.lower().endswith('.png')
        ])
        
        if not img_paths:
            raise Exception("No image files found in the output directory")
            
        print(f"Found {len(img_paths)} image files to process")
        
        video_path = os.path.join(CONVERTED_DIR, os.path.basename(pptx_path).replace(".pptx", ".mp4"))
        print(f"Target output path: {video_path}")
        
        # Check if ffmpeg is available
        if not ffmpeg_available:
            print("FFmpeg not available, creating animated GIF instead")
            result_path = create_fallback_slideshow(img_paths, video_path)
            
            # Inform the user about the fallback
            print("Created an animated GIF instead of MP4 due to missing FFmpeg")
            
            # Verify file exists
            if not os.path.exists(result_path):
                raise Exception("Failed to create slideshow")
                
            return result_path
        
        # FFmpeg is available, create video
        print("Creating video with FFmpeg")
        clip = ImageSequenceClip(img_paths, fps=1)
        clip.write_videofile(video_path, codec="libx264", audio=False)
        
        # Verify video was created
        if not os.path.exists(video_path):
            raise Exception("Video file was not created")
            
        print(f"Video conversion completed successfully")
        return video_path
        
    except Exception as e:
        print(f"Error in PowerPoint to Video conversion: {str(e)}")
        traceback_info = traceback.format_exc()
        print(f"Traceback: {traceback_info}")
        
        # If this was a FFmpeg error, suggest installation
        if "FFmpeg is required" in str(e):
            message = (
                "PowerPoint to Video conversion requires FFmpeg. "
                "Please install FFmpeg and add it to your system PATH. "
                "Installation instructions:\n"
                "- Windows: Download from https://ffmpeg.org/download.html and add to PATH\n"
                "- macOS: Run 'brew install ffmpeg' if you have Homebrew\n"
                "- Linux: Run 'sudo apt install ffmpeg' (Ubuntu/Debian) or equivalent for your distro"
            )
            raise Exception(message)
        
        raise Exception(f"PowerPoint to Video conversion failed: {str(e)}")

# ✅ Text to PDF
def convert_text_to_pdf(txt_path):
    try:
        print(f"Starting Text to PDF conversion: {txt_path}")
        pdf_path = os.path.join(CONVERTED_DIR, os.path.basename(txt_path).replace(".txt", ".pdf"))
        print(f"Output will be saved to: {pdf_path}")
        
        # Try with pdfkit first
        try:
            print("Attempting conversion with pdfkit/wkhtmltopdf...")
            pdfkit.from_file(txt_path, pdf_path)
            print(f"PDF created successfully with pdfkit: {pdf_path}")
            return pdf_path
        except Exception as e:
            print(f"pdfkit conversion failed: {str(e)}")
            print("Trying alternative method...")
        
        # Fallback to ReportLab
        try:
            print("Using ReportLab as fallback method...")
            
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.lib.units import inch
            
            # Read text file
            with open(txt_path, 'r', encoding='utf-8') as f:
                text_content = f.read()
            
            # Create PDF
            document = SimpleDocTemplate(pdf_path, pagesize=letter)
            styles = getSampleStyleSheet()
            
            # Create content elements
            elements = []
            
            # Split text by lines
            lines = text_content.split('\n')
            
            # Process the lines
            current_para = []
            for line in lines:
                # Empty line indicates paragraph break
                if not line.strip():
                    if current_para:
                        para_text = ' '.join(current_para)
                        elements.append(Paragraph(para_text, styles['Normal']))
                        elements.append(Spacer(1, 0.1*inch))
                        current_para = []
                else:
                    current_para.append(line)
            
            # Add the last paragraph if any
            if current_para:
                para_text = ' '.join(current_para)
                elements.append(Paragraph(para_text, styles['Normal']))
            
            # Build the PDF
            document.build(elements)
            
            print(f"PDF created successfully with ReportLab: {pdf_path}")
            return pdf_path
            
        except Exception as e:
            print(f"ReportLab conversion failed: {str(e)}")
            raise Exception(f"All PDF conversion methods failed")
    
    except Exception as e:
        print(f"Error in Text to PDF conversion: {str(e)}")
        traceback_info = traceback.format_exc()
        print(f"Traceback: {traceback_info}")
        raise Exception(f"Text to PDF conversion failed: {str(e)}")

# ✅ Text to Word
def convert_text_to_docx(txt_path):
    docx_path = os.path.join(CONVERTED_DIR, os.path.basename(txt_path).replace(".txt", ".docx"))
    doc = Document()
    with open(txt_path, "r", encoding="utf-8") as f:
        doc.add_paragraph(f.read())
    doc.save(docx_path)
    return docx_path

# ✅ Text to Markdown
def convert_text_to_markdown(txt_path):
    md_path = os.path.join(CONVERTED_DIR, os.path.basename(txt_path).replace(".txt", ".md"))
    with open(txt_path, "r", encoding="utf-8") as f:
        text = f.read()
    with open(md_path, "w", encoding="utf-8") as f:
        f.write(text)
    return md_path

# ✅ HTML to PDF (Using WeasyPrint)
async def convert_html_to_pdf(html_path):
    pdf_path = os.path.join(CONVERTED_DIR, os.path.basename(html_path).replace(".html", "_converted.pdf"))
    HTML(html_path).write_pdf(pdf_path)
    return pdf_path

# ✅ HTML to DOCX
def convert_html_to_docx(html_path):
    docx_path = os.path.join(CONVERTED_DIR, os.path.basename(html_path).replace(".html", ".docx"))
    with open(html_path, "r", encoding="utf-8") as f:
        html_content = f.read()
    
    doc = Document()
    doc.add_paragraph(html_content)
    doc.save(docx_path)
    return docx_path

# ✅ Image to PDF (Using ReportLab for more reliable conversion)
async def convert_image_to_pdf(image_path):
    try:
        print(f"Starting Image to PDF conversion: {image_path}")
        pdf_path = os.path.join(CONVERTED_DIR, os.path.splitext(os.path.basename(image_path))[0] + "_converted.pdf")
        print(f"Output will be saved to: {pdf_path}")
        
        # Check if image exists
        if not os.path.exists(image_path):
            error_msg = f"Image file does not exist: {image_path}"
            print(error_msg)
            raise Exception(error_msg)
        
        # Open the image
        try:
            image = Image.open(image_path)
            print(f"Image opened successfully: {image.format}, size: {image.size}, mode: {image.mode}")
        except Exception as e:
            error_msg = f"Failed to open image: {str(e)}"
            print(error_msg)
            raise Exception(error_msg)
        
        # Convert to RGB (required for PDF)
        if image.mode != 'RGB':
            try:
                image = image.convert('RGB')
                print(f"Converted image to RGB mode")
            except Exception as e:
                print(f"Warning: Failed to convert image to RGB mode: {str(e)}")
                # Try to continue anyway
        
        # Try saving with reportlab (more reliable)
        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            
            print("Using ReportLab for PDF creation")
            
            # Use reportlab for PDF creation (often more reliable)
            c = canvas.Canvas(pdf_path, pagesize=letter)
            width, height = letter
            
            # Calculate image dimensions to fit the page while maintaining aspect ratio
            img_width, img_height = image.size
            aspect = img_height / float(img_width)
            
            # Keep the width but adjust the height according to aspect ratio
            display_width = width * 0.9  # 90% of page width
            display_height = display_width * aspect
            
            # Center on page
            x = (width - display_width) / 2
            y = (height - display_height) / 2
            
            # Save the image temporarily
            temp_img_path = os.path.join(CONVERTED_DIR, "temp_img.png")
            image.save(temp_img_path)
            
            # Add the image to the canvas
            c.drawImage(temp_img_path, x, y, width=display_width, height=display_height)
            c.save()
            
            # Clean up temp file
            if os.path.exists(temp_img_path):
                os.remove(temp_img_path)
                
            reportlab_success = True
        except Exception as e:
            print(f"ReportLab approach failed: {str(e)}. Falling back to PIL...")
            reportlab_success = False
        
        # If reportlab didn't work, try with PIL
        if not reportlab_success:
            try:
                # Use PIL for PDF creation (fallback)
                image.save(
                    pdf_path, 
                    "PDF", 
                    resolution=300,
                    quality=100,
                    optimize=True
                )
            except Exception as e:
                error_msg = f"Failed to convert image to PDF with PIL: {str(e)}"
                print(error_msg)
                
                # Last resort: try using img2pdf if available
                try:
                    import img2pdf
                    print("Trying img2pdf as last resort")
                    
                    # Convert image to RGB and save as temporary file
                    temp_rgb_path = os.path.join(CONVERTED_DIR, "temp_rgb.png")
                    image.convert('RGB').save(temp_rgb_path)
                    
                    # Convert to PDF using img2pdf
                    with open(pdf_path, "wb") as f:
                        f.write(img2pdf.convert(temp_rgb_path))
                    
                    # Clean up
                    if os.path.exists(temp_rgb_path):
                        os.remove(temp_rgb_path)
                except ImportError:
                    print("img2pdf not available")
                    raise Exception("All PDF conversion methods failed")
                except Exception as e:
                    print(f"img2pdf approach failed: {str(e)}")
                    raise Exception("All PDF conversion methods failed")
        
        # Verify the file was created
        if os.path.exists(pdf_path):
            file_size = os.path.getsize(pdf_path)
            print(f"PDF created successfully, size: {file_size} bytes")
            
            if file_size == 0:
                raise Exception("PDF file was created but is empty")
                
            return pdf_path
        else:
            raise Exception("PDF file was not created")
            
    except Exception as e:
        print(f"Error in Image to PDF conversion: {str(e)}")
        traceback_info = traceback.format_exc()
        print(f"Traceback: {traceback_info}")
        raise Exception(f"Image to PDF conversion failed: {str(e)}")

# ✅ CSV to JSON
def convert_csv_to_json(csv_path):
    json_path = os.path.join(CONVERTED_DIR, os.path.basename(csv_path).replace(".csv", ".json"))
    df = pd.read_csv(csv_path)
    df.to_json(json_path, orient="records", indent=4)
    return json_path

# ✅ CSV to Excel (XLSX)
def convert_csv_to_xlsx(csv_path):
    xlsx_path = os.path.join(CONVERTED_DIR, os.path.basename(csv_path).replace(".csv", ".xlsx"))
    df = pd.read_csv(csv_path)
    df.to_excel(xlsx_path, index=False)
    return xlsx_path

# ✅ Markdown to HTML
def convert_markdown_to_html(md_path):
    html_path = os.path.join(CONVERTED_DIR, os.path.basename(md_path).replace(".md", ".html"))
    with open(md_path, "r", encoding="utf-8") as f:
        md_content = f.read()
    html_content = markdown.markdown(md_content)
    with open(html_path, "w", encoding="utf-8") as f:
        f.write(html_content)
    return html_path

# ✅ Markdown to DOCX
def convert_markdown_to_docx(md_path):
    docx_path = os.path.join(CONVERTED_DIR, os.path.basename(md_path).replace(".md", ".docx"))
    with open(md_path, "r", encoding="utf-8") as f:
        md_content = f.read()
    doc = Document()
    for line in md_content.split("\n"):
        doc.add_paragraph(line)
    doc.save(docx_path)
    return docx_path

# Optimized image format conversion
async def convert_image_format(image_path, target_format):
    """Convert an image to a different format with simple HEIC support"""
    try:
        print(f"Starting image format conversion: {image_path} to {target_format}")
        
        # Validate target format
        supported_formats = {'jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff', 'webp', 'ico', 'heic'}
        target_format = target_format.lower()
        if target_format not in supported_formats:
            raise ValueError(f"Unsupported target format: {target_format}. Supported formats: {', '.join(supported_formats)}")
        
        # Create output path
        output_filename = f"{os.path.splitext(os.path.basename(image_path))[0]}_converted.{target_format}"
        output_path = os.path.join(CONVERTED_DIR, output_filename)
        print(f"Output will be saved as: {output_path}")
        
        # Check if this is a HEIC file
        is_heic_input = image_path.lower().endswith('.heic')
        is_heic_output = target_format.lower() == 'heic'
        
        # HEIC to standard format conversion
        if is_heic_input and not is_heic_output:
            print("Converting from HEIC to standard format")
            try:
                # Try using FFmpeg first for HEIC conversion (most reliable)
                ffmpeg_path = None
                try:
                    import imageio_ffmpeg
                    ffmpeg_path = imageio_ffmpeg.get_ffmpeg_exe()
                except:
                    pass
                
                if not ffmpeg_path:
                    local_ffmpeg = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ffmpeg", "ffmpeg.exe")
                    if os.path.exists(local_ffmpeg):
                        ffmpeg_path = local_ffmpeg
                
                if ffmpeg_path and os.path.exists(ffmpeg_path):
                    print(f"Using FFmpeg for HEIC conversion: {ffmpeg_path}")
                    cmd = [
                        ffmpeg_path, "-y", "-i", image_path,
                        output_path
                    ]
                    result = subprocess.run(
                        cmd,
                        stdout=subprocess.PIPE,
                        stderr=subprocess.PIPE,
                        text=True
                    )
                    
                    if result.returncode == 0 and os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                        print(f"FFmpeg HEIC conversion successful: {output_path}")
                        return output_path
                    else:
                        print(f"FFmpeg conversion failed, trying pillow-heif: {result.stderr}")
                else:
                    print("FFmpeg not found, trying pillow-heif")
                
                # If FFmpeg fails or is not available, try pillow-heif
                from pillow_heif import register_heif_opener
                register_heif_opener()
                print("Registered HEIF opener")
                
                # Use PIL to open and convert
                img = Image.open(image_path)
                if img.mode not in ['RGB', 'RGBA'] and target_format not in ['bmp', 'gif']:
                    img = img.convert('RGB')
                
                # Save options for different formats
                save_opts = {}
                if target_format == 'jpg' or target_format == 'jpeg':
                    save_opts = {'quality': 95, 'optimize': True}
                elif target_format == 'png':
                    save_opts = {'optimize': True}
                elif target_format == 'webp':
                    save_opts = {'quality': 90}
                
                img.save(output_path, format=target_format.upper(), **save_opts)
                print(f"HEIC conversion with PIL successful: {output_path}")
                return output_path
                
            except Exception as e:
                print(f"Error in HEIC conversion: {str(e)}")
                traceback.print_exc()
                raise Exception(f"HEIC conversion failed: {str(e)}")
        
        # Standard format to HEIC conversion
        elif is_heic_output:
            print("Converting to HEIC format")
            try:
                # Import and use pillow-heif
                from pillow_heif import register_heif_opener, from_pillow
                register_heif_opener()
                
                # Open with PIL
                img = Image.open(image_path)
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                
                # Convert to HEIC
                heif_file = from_pillow(img)
                heif_file.save(output_path)
                
                print(f"HEIC conversion successful: {output_path}")
                return output_path
            except Exception as e:
                print(f"Error converting to HEIC: {str(e)}")
                traceback.print_exc()
                
                # Fallback to PNG
                fallback_path = os.path.join(CONVERTED_DIR, f"{os.path.splitext(os.path.basename(image_path))[0]}_converted.png")
                try:
                    img.save(fallback_path, format='PNG')
                    print(f"Fallback to PNG: {fallback_path}")
                    return fallback_path
                except:
                    raise Exception(f"HEIC conversion failed: {str(e)}")
        
        # Standard format conversion (non-HEIC)
        else:
            print("Standard format conversion")
            try:
                # Open image with PIL
                img = Image.open(image_path)
                
                # Convert color mode if needed
                if target_format in ['jpg', 'jpeg'] and img.mode not in ['RGB', 'RGBA']:
                    img = img.convert('RGB')
                
                # Format-specific options
                save_opts = {}
                if target_format == 'jpg' or target_format == 'jpeg':
                    save_opts = {'quality': 95, 'optimize': True}
                elif target_format == 'png':
                    save_opts = {'optimize': True}
                elif target_format == 'webp':
                    save_opts = {'quality': 90}
                elif target_format == 'ico':
                    img = img.resize((32, 32))
                
                # Save the image
                img.save(output_path, format=target_format.upper(), **save_opts)
                print(f"Standard conversion successful: {output_path}")
                return output_path
            except Exception as e:
                print(f"Error in standard conversion: {str(e)}")
                traceback.print_exc()
                raise Exception(f"Image conversion failed: {str(e)}")
    
    except Exception as e:
        print(f"Error in convert_image_format: {str(e)}")
        traceback.print_exc()
        raise Exception(f"Image conversion failed: {str(e)}")

# ✅ Image Compression
async def compress_image(image_path, quality=85):
    try:
        print(f"Starting image compression: {image_path} with quality {quality}%")
        
        # Validate quality parameter
        quality = int(quality)
        if quality < 1 or quality > 100:
            quality = 85  # Use default if outside valid range
            print(f"Invalid quality value, using default quality: {quality}%")
        
        # Get original file info
        original_size = os.path.getsize(image_path)
        original_format = os.path.splitext(image_path)[1].lower()
        
        # Determine output path - preserve original format
        output_filename = f"{os.path.splitext(os.path.basename(image_path))[0]}_compressed{original_format}"
        output_path = os.path.join(CONVERTED_DIR, output_filename)
        print(f"Compressed image will be saved as: {output_path}")
        
        # Open and process the image
        with Image.open(image_path) as img:
            print(f"Original image: {img.format}, mode: {img.mode}, size: {img.width}x{img.height}, file size: {original_size/1024:.1f} KB")
            
            # Preserve original mode (especially important for PNG transparency)
            original_mode = img.mode
            
            # Optimize based on image format
            if original_format == '.png':
                # For PNG, use optimize and specific compression level
                img.save(output_path, 
                        format='PNG',
                        optimize=True,
                        quality=quality,
                        compress_level=8)  # Higher compression level (1-9)
            elif original_format in ['.jpg', '.jpeg']:
                # For JPEG, focus on quality
                if original_mode == 'RGBA':
                    # Convert to RGB if needed
                    img = img.convert('RGB')
                img.save(output_path,
                        format='JPEG',
                        quality=quality,
                        optimize=True)
            else:
                # For other formats, try to preserve original format
                img.save(output_path,
                        format=img.format if img.format else 'PNG',
                        quality=quality,
                        optimize=True)
            
        # Check compression results
        if os.path.exists(output_path):
            compressed_size = os.path.getsize(output_path)
            reduction = (1 - compressed_size/original_size) * 100
            
            if compressed_size < original_size:
                print(f"Compression successful: {compressed_size/1024:.1f} KB ({reduction:.1f}% reduction)")
                return output_path
            else:
                print("Compressed file is larger than original, attempting alternative compression...")
                
                # Try alternative compression method
                with Image.open(image_path) as img:
                    if original_format == '.png':
                        # Try with different compression parameters
                        alternative_path = os.path.join(CONVERTED_DIR, f"{os.path.splitext(os.path.basename(image_path))[0]}_compressed_alt.png")
                        img.save(alternative_path,
                                format='PNG',
                                optimize=True,
                                compress_level=9)  # Maximum compression
                        
                        alt_size = os.path.getsize(alternative_path)
                        if alt_size < original_size:
                            print(f"Alternative compression successful: {alt_size/1024:.1f} KB ({(1 - alt_size/original_size) * 100:.1f}% reduction)")
                            return alternative_path
                
                print("Could not achieve better compression. Returning original file.")
                return image_path
        else:
            raise Exception("Failed to save compressed image")
            
    except Exception as e:
        print(f"Error in image compression: {str(e)}")
        traceback_info = traceback.format_exc()
        print(f"Traceback: {traceback_info}")
        raise Exception(f"Image compression failed: {str(e)}")

# Add this function near other conversion functions
async def extract_audio_from_video(video_path, output_format="mp3", task_id=None):
    """Extract audio from a video file and save in specified format"""
    try:
        print(f"Starting audio extraction from video: {video_path}")
        
        # Update progress if task_id provided
        if task_id:
            conversion_progress[task_id] = {"progress": 10, "status": "Starting audio extraction"}
        
        # Create output path based on input filename
        output_path = os.path.join(CONVERTED_DIR, os.path.splitext(os.path.basename(video_path))[0] + f".{output_format}")
        
        # Check if ffmpeg is available
        if not ffmpeg_available:
            error_msg = "FFmpeg is required for audio extraction but not available"
            print(error_msg)
            if task_id:
                conversion_progress[task_id] = {"progress": 100, "status": f"Error: {error_msg}"}
            raise Exception(error_msg)
        
        # Get FFmpeg path - try the detected one or the local one
        ffmpeg_executable = ffmpeg_path
        if not ffmpeg_executable or not os.path.exists(ffmpeg_executable):
            local_ffmpeg = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ffmpeg", "ffmpeg.exe")
            if os.path.exists(local_ffmpeg):
                ffmpeg_executable = local_ffmpeg
        
        if not ffmpeg_executable or not os.path.exists(ffmpeg_executable):
            # Try to get it from imageio
            try:
                import imageio_ffmpeg
                ffmpeg_executable = imageio_ffmpeg.get_ffmpeg_exe()
            except Exception:
                pass
        
        if not ffmpeg_executable or not os.path.exists(ffmpeg_executable):
            error_msg = "Cannot locate FFmpeg executable"
            print(error_msg)
            if task_id:
                conversion_progress[task_id] = {"progress": 100, "status": f"Error: {error_msg}"}
            raise Exception(error_msg)
            
        if task_id:
            conversion_progress[task_id] = {"progress": 30, "status": "Preparing conversion"}
        
        # For the basic formats (MP3, WAV, AAC, OGG) that worked before, we can still try MoviePy first
        if output_format.lower() in ["mp3", "wav", "aac", "ogg"]:
            try:
                # Import moviepy's VideoFileClip
                try:
                    # Try direct import first (newer versions of moviepy)
                    from moviepy.video.io.VideoFileClip import VideoFileClip
                except ImportError:
                    # Fall back to editor import (older versions)
                    from moviepy.editor import VideoFileClip
                
                if task_id:
                    conversion_progress[task_id] = {"progress": 40, "status": "Loading video file"}
                
                # Load the video file
                video = VideoFileClip(video_path)
                
                if task_id:
                    conversion_progress[task_id] = {"progress": 50, "status": "Extracting audio track"}
                
                # Extract audio
                audio = video.audio
                
                if audio is None:
                    error_msg = "No audio track found in the video file"
                    print(error_msg)
                    if task_id:
                        conversion_progress[task_id] = {"progress": 100, "status": f"Error: {error_msg}"}
                    raise Exception(error_msg)
                
                if task_id:
                    conversion_progress[task_id] = {"progress": 70, "status": f"Converting to {output_format}"}
                
                # Write audio to file using MoviePy
                if output_format.lower() == "mp3":
                    audio.write_audiofile(output_path, codec='mp3')
                elif output_format.lower() == "wav":
                    audio.write_audiofile(output_path, codec='pcm_s16le')
                elif output_format.lower() == "aac":
                    audio.write_audiofile(output_path, codec='aac')
                elif output_format.lower() == "ogg":
                    audio.write_audiofile(output_path, codec='libvorbis')
                else:
                    # Default to mp3
                    audio.write_audiofile(output_path, codec='mp3')
                
                # Close the clips to release resources
                audio.close()
                video.close()
                
                # Verify file was created
                if not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
                    raise Exception("MoviePy did not create the output file successfully")
                
                if task_id:
                    conversion_progress[task_id] = {"progress": 100, "status": "Audio extraction complete", "file_path": output_path, "file_name": os.path.basename(output_path)}
                
                print(f"Audio extraction complete: {output_path}")
                return output_path
                
            except Exception as e:
                print(f"MoviePy approach failed: {str(e)}. Falling back to direct FFmpeg.")
                # Fall back to direct FFmpeg conversion
                return direct_ffmpeg_conversion(video_path, output_format, ffmpeg_executable, task_id)
        else:
            # For all other formats, go straight to direct FFmpeg
            return direct_ffmpeg_conversion(video_path, output_format, ffmpeg_executable, task_id)
    
    except Exception as e:
        print(f"Error extracting audio: {str(e)}")
        traceback_info = traceback.format_exc()
        print(f"Traceback: {traceback_info}")
        
        if task_id:
            conversion_progress[task_id] = {"progress": 100, "status": f"Error: {str(e)}"}
        
        raise Exception(f"Audio extraction failed: {str(e)}")

# Define the mapping of file extensions and target formats to conversion functions
conversion_functions = {
    "pdf": {
        "docx": convert_pdf_to_docx,
        "xlsx": convert_pdf_to_excel,
        "pptx": convert_pdf_to_pptx,
        "txt": convert_pdf_to_text,
        "html": convert_pdf_to_html
    },
    "docx": {
        "pdf": convert_docx_to_pdf,
        "html": convert_docx_to_html,
        "md": convert_docx_to_markdown
    },
    "xlsx": {
        "csv": convert_xlsx_to_csv,
        "json": convert_xlsx_to_json,
        "xml": convert_xlsx_to_xml
    },
    "pptx": {
        "pdf": convert_pptx_to_pdf,
        "images": convert_pptx_to_images,
        "video": convert_pptx_to_video
    },
    "txt": {
        "pdf": convert_text_to_pdf,
        "docx": convert_text_to_docx,
        "md": convert_text_to_markdown
    },
    "html": {
        "pdf": convert_html_to_pdf,
        "docx": convert_html_to_docx
    },
    "mp4": {
        "mp3": lambda path: extract_audio_from_video(path, "mp3"),
        "wav": lambda path: extract_audio_from_video(path, "wav"),
        "aac": lambda path: extract_audio_from_video(path, "aac"),
        "ogg": lambda path: extract_audio_from_video(path, "ogg")
    },
    "avi": {
        "mp3": lambda path: extract_audio_from_video(path, "mp3"),
        "wav": lambda path: extract_audio_from_video(path, "wav"),
        "aac": lambda path: extract_audio_from_video(path, "aac"),
        "ogg": lambda path: extract_audio_from_video(path, "ogg")
    },
    "mov": {
        "mp3": lambda path: extract_audio_from_video(path, "mp3"),
        "wav": lambda path: extract_audio_from_video(path, "wav"),
        "aac": lambda path: extract_audio_from_video(path, "aac"),
        "ogg": lambda path: extract_audio_from_video(path, "ogg")
    },
    "mkv": {
        "mp3": lambda path: extract_audio_from_video(path, "mp3"),
        "wav": lambda path: extract_audio_from_video(path, "wav"),
        "aac": lambda path: extract_audio_from_video(path, "aac"),
        "ogg": lambda path: extract_audio_from_video(path, "ogg")
    },
    "webm": {
        "mp3": lambda path: extract_audio_from_video(path, "mp3"),
        "wav": lambda path: extract_audio_from_video(path, "wav"),
        "aac": lambda path: extract_audio_from_video(path, "aac"),
        "ogg": lambda path: extract_audio_from_video(path, "ogg")
    },
    "flv": {
        "mp3": lambda path: extract_audio_from_video(path, "mp3"),
        "wav": lambda path: extract_audio_from_video(path, "wav"),
        "aac": lambda path: extract_audio_from_video(path, "aac"),
        "ogg": lambda path: extract_audio_from_video(path, "ogg")
    },
    "wmv": {
        "mp3": lambda path: extract_audio_from_video(path, "mp3"),
        "wav": lambda path: extract_audio_from_video(path, "wav"),
        "aac": lambda path: extract_audio_from_video(path, "aac"),
        "ogg": lambda path: extract_audio_from_video(path, "ogg"),
        "flac": lambda path: extract_audio_from_video(path, "flac"),
        "opus": lambda path: extract_audio_from_video(path, "opus"),
        "ac3": lambda path: extract_audio_from_video(path, "ac3"),
        "m4a": lambda path: extract_audio_from_video(path, "m4a"),
        "alac": lambda path: extract_audio_from_video(path, "alac")
    },
    "jpg": {
        "pdf": convert_image_to_pdf,
        "text (ocr)": extract_text_from_image,
        "png": lambda path: convert_image_format(path, "png"),
        "webp": lambda path: convert_image_format(path, "webp"),
        "gif": lambda path: convert_image_format(path, "gif"),
        "bmp": lambda path: convert_image_format(path, "bmp"),
        "tiff": lambda path: convert_image_format(path, "tiff"),
        "ico": lambda path: convert_image_format(path, "ico"),
        "jpeg": lambda path: convert_image_format(path, "jpeg"),
        "compressed": compress_image,
        "heic": lambda path: convert_image_format(path, "heic")
    },
    "jpeg": {
        "pdf": convert_image_to_pdf,
        "text (ocr)": extract_text_from_image,
        "png": lambda path: convert_image_format(path, "png"),
        "webp": lambda path: convert_image_format(path, "webp"),
        "gif": lambda path: convert_image_format(path, "gif"),
        "bmp": lambda path: convert_image_format(path, "bmp"),
        "tiff": lambda path: convert_image_format(path, "tiff"),
        "ico": lambda path: convert_image_format(path, "ico"),
        "jpg": lambda path: convert_image_format(path, "jpg"),
        "compressed": compress_image,
        "heic": lambda path: convert_image_format(path, "heic")
    },
    "png": {
        "pdf": convert_image_to_pdf,
        "text (ocr)": extract_text_from_image,
        "jpg": lambda path: convert_image_format(path, "jpg"),
        "jpeg": lambda path: convert_image_format(path, "jpeg"),
        "webp": lambda path: convert_image_format(path, "webp"),
        "gif": lambda path: convert_image_format(path, "gif"),
        "bmp": lambda path: convert_image_format(path, "bmp"),
        "tiff": lambda path: convert_image_format(path, "tiff"),
        "ico": lambda path: convert_image_format(path, "ico"),
        "compressed": compress_image,
        "heic": lambda path: convert_image_format(path, "heic")
    },
    "webp": {
        "pdf": convert_image_to_pdf,
        "text (ocr)": extract_text_from_image,
        "png": lambda path: convert_image_format(path, "png"),
        "jpg": lambda path: convert_image_format(path, "jpg"),
        "jpeg": lambda path: convert_image_format(path, "jpeg"),
        "gif": lambda path: convert_image_format(path, "gif"),
        "bmp": lambda path: convert_image_format(path, "bmp"),
        "tiff": lambda path: convert_image_format(path, "tiff"),
        "ico": lambda path: convert_image_format(path, "ico"),
        "compressed": compress_image,
        "heic": lambda path: convert_image_format(path, "heic")
    },
    "gif": {
        "pdf": convert_image_to_pdf,
        "png": lambda path: convert_image_format(path, "png"),
        "jpg": lambda path: convert_image_format(path, "jpg"),
        "jpeg": lambda path: convert_image_format(path, "jpeg"),
        "webp": lambda path: convert_image_format(path, "webp"),
        "bmp": lambda path: convert_image_format(path, "bmp"),
        "tiff": lambda path: convert_image_format(path, "tiff"),
        "ico": lambda path: convert_image_format(path, "ico"),
        "compressed": compress_image,
        "heic": lambda path: convert_image_format(path, "heic")
    },
    "bmp": {
        "pdf": convert_image_to_pdf,
        "text (ocr)": extract_text_from_image,
        "png": lambda path: convert_image_format(path, "png"),
        "jpg": lambda path: convert_image_format(path, "jpg"),
        "jpeg": lambda path: convert_image_format(path, "jpeg"),
        "webp": lambda path: convert_image_format(path, "webp"),
        "gif": lambda path: convert_image_format(path, "gif"),
        "tiff": lambda path: convert_image_format(path, "tiff"),
        "ico": lambda path: convert_image_format(path, "ico"),
        "compressed": compress_image,
        "heic": lambda path: convert_image_format(path, "heic")
    },
    "tiff": {
        "pdf": convert_image_to_pdf,
        "text (ocr)": extract_text_from_image,
        "png": lambda path: convert_image_format(path, "png"),
        "jpg": lambda path: convert_image_format(path, "jpg"),
        "jpeg": lambda path: convert_image_format(path, "jpeg"),
        "webp": lambda path: convert_image_format(path, "webp"),
        "gif": lambda path: convert_image_format(path, "gif"),
        "bmp": lambda path: convert_image_format(path, "bmp"),
        "ico": lambda path: convert_image_format(path, "ico"),
        "compressed": compress_image,
        "heic": lambda path: convert_image_format(path, "heic")
    },
    "ico": {
        "png": lambda path: convert_image_format(path, "png"),
        "jpg": lambda path: convert_image_format(path, "jpg"),
        "jpeg": lambda path: convert_image_format(path, "jpeg"),
        "webp": lambda path: convert_image_format(path, "webp"),
        "gif": lambda path: convert_image_format(path, "gif"),
        "bmp": lambda path: convert_image_format(path, "bmp"),
        "tiff": lambda path: convert_image_format(path, "tiff"),
        "compressed": compress_image,
        "heic": lambda path: convert_image_format(path, "heic")
    },
    "csv": {
        "json": convert_csv_to_json, 
        "xlsx": convert_csv_to_xlsx
    },
    "json": {
        "xml": convert_json_to_xml
    },
    "xml": {
        "json": convert_xml_to_json
    },
    "md": {
        "html": convert_markdown_to_html, 
        "docx": convert_markdown_to_docx
    },
    # Video format conversions to audio
    "mp4": {
        "mp3": lambda path: extract_audio_from_video(path, "mp3"),
        "wav": lambda path: extract_audio_from_video(path, "wav"),
        "aac": lambda path: extract_audio_from_video(path, "aac"),
        "ogg": lambda path: extract_audio_from_video(path, "ogg"),
        "flac": lambda path: extract_audio_from_video(path, "flac"),
        "opus": lambda path: extract_audio_from_video(path, "opus"),
        "ac3": lambda path: extract_audio_from_video(path, "ac3"),
        "m4a": lambda path: extract_audio_from_video(path, "m4a"),
        "alac": lambda path: extract_audio_from_video(path, "alac")
    },
    "avi": {
        "mp3": lambda path: extract_audio_from_video(path, "mp3"),
        "wav": lambda path: extract_audio_from_video(path, "wav"),
        "aac": lambda path: extract_audio_from_video(path, "aac"),
        "ogg": lambda path: extract_audio_from_video(path, "ogg"),
        "flac": lambda path: extract_audio_from_video(path, "flac"),
        "opus": lambda path: extract_audio_from_video(path, "opus"),
        "ac3": lambda path: extract_audio_from_video(path, "ac3"),
        "m4a": lambda path: extract_audio_from_video(path, "m4a"),
        "alac": lambda path: extract_audio_from_video(path, "alac")
    },
    "mov": {
        "mp3": lambda path: extract_audio_from_video(path, "mp3"),
        "wav": lambda path: extract_audio_from_video(path, "wav"),
        "aac": lambda path: extract_audio_from_video(path, "aac"),
        "ogg": lambda path: extract_audio_from_video(path, "ogg"),
        "flac": lambda path: extract_audio_from_video(path, "flac"),
        "opus": lambda path: extract_audio_from_video(path, "opus"),
        "ac3": lambda path: extract_audio_from_video(path, "ac3"),
        "m4a": lambda path: extract_audio_from_video(path, "m4a"),
        "alac": lambda path: extract_audio_from_video(path, "alac")
    },
    "mkv": {
        "mp3": lambda path: extract_audio_from_video(path, "mp3"),
        "wav": lambda path: extract_audio_from_video(path, "wav"),
        "aac": lambda path: extract_audio_from_video(path, "aac"),
        "ogg": lambda path: extract_audio_from_video(path, "ogg"),
        "flac": lambda path: extract_audio_from_video(path, "flac"),
        "opus": lambda path: extract_audio_from_video(path, "opus"),
        "ac3": lambda path: extract_audio_from_video(path, "ac3"),
        "m4a": lambda path: extract_audio_from_video(path, "m4a"),
        "alac": lambda path: extract_audio_from_video(path, "alac")
    },
    "webm": {
        "mp3": lambda path: extract_audio_from_video(path, "mp3"),
        "wav": lambda path: extract_audio_from_video(path, "wav"),
        "aac": lambda path: extract_audio_from_video(path, "aac"),
        "ogg": lambda path: extract_audio_from_video(path, "ogg"),
        "flac": lambda path: extract_audio_from_video(path, "flac"),
        "opus": lambda path: extract_audio_from_video(path, "opus"),
        "ac3": lambda path: extract_audio_from_video(path, "ac3"),
        "m4a": lambda path: extract_audio_from_video(path, "m4a"),
        "alac": lambda path: extract_audio_from_video(path, "alac")
    },
    "flv": {
        "mp3": lambda path: extract_audio_from_video(path, "mp3"),
        "wav": lambda path: extract_audio_from_video(path, "wav"),
        "aac": lambda path: extract_audio_from_video(path, "aac"),
        "ogg": lambda path: extract_audio_from_video(path, "ogg"),
        "flac": lambda path: extract_audio_from_video(path, "flac"),
        "opus": lambda path: extract_audio_from_video(path, "opus"),
        "ac3": lambda path: extract_audio_from_video(path, "ac3"),
        "m4a": lambda path: extract_audio_from_video(path, "m4a"),
        "alac": lambda path: extract_audio_from_video(path, "alac")
    },
    "wmv": {
        "mp3": lambda path: extract_audio_from_video(path, "mp3"),
        "wav": lambda path: extract_audio_from_video(path, "wav"),
        "aac": lambda path: extract_audio_from_video(path, "aac"),
        "ogg": lambda path: extract_audio_from_video(path, "ogg"),
        "flac": lambda path: extract_audio_from_video(path, "flac"),
        "opus": lambda path: extract_audio_from_video(path, "opus"),
        "ac3": lambda path: extract_audio_from_video(path, "ac3"),
        "m4a": lambda path: extract_audio_from_video(path, "m4a"),
        "alac": lambda path: extract_audio_from_video(path, "alac")
    },
    "mp3": {
        "wav": lambda path: convert_audio_format(path, "wav"),
        "aac": lambda path: convert_audio_format(path, "aac"),
        "ogg": lambda path: convert_audio_format(path, "ogg"),
        "flac": lambda path: convert_audio_format(path, "flac"),
        "opus": lambda path: convert_audio_format(path, "opus"),
        "ac3": lambda path: convert_audio_format(path, "ac3"),
        "m4a": lambda path: convert_audio_format(path, "m4a"),
        "alac": lambda path: convert_audio_format(path, "alac")
    },
    "wav": {
        "mp3": lambda path: convert_audio_format(path, "mp3"),
        "aac": lambda path: convert_audio_format(path, "aac"),
        "ogg": lambda path: convert_audio_format(path, "ogg"),
        "flac": lambda path: convert_audio_format(path, "flac"),
        "opus": lambda path: convert_audio_format(path, "opus"),
        "ac3": lambda path: convert_audio_format(path, "ac3"),
        "m4a": lambda path: convert_audio_format(path, "m4a"),
        "alac": lambda path: convert_audio_format(path, "alac")
    },
    "aac": {
        "mp3": lambda path: convert_audio_format(path, "mp3"),
        "wav": lambda path: convert_audio_format(path, "wav"),
        "ogg": lambda path: convert_audio_format(path, "ogg"),
        "flac": lambda path: convert_audio_format(path, "flac"),
        "opus": lambda path: convert_audio_format(path, "opus"),
        "ac3": lambda path: convert_audio_format(path, "ac3"),
        "m4a": lambda path: convert_audio_format(path, "m4a"),
        "alac": lambda path: convert_audio_format(path, "alac")
    },
    "ogg": {
        "mp3": lambda path: convert_audio_format(path, "mp3"),
        "wav": lambda path: convert_audio_format(path, "wav"),
        "aac": lambda path: convert_audio_format(path, "aac"),
        "flac": lambda path: convert_audio_format(path, "flac"),
        "opus": lambda path: convert_audio_format(path, "opus"),
        "ac3": lambda path: convert_audio_format(path, "ac3"),
        "m4a": lambda path: convert_audio_format(path, "m4a"),
        "alac": lambda path: convert_audio_format(path, "alac")
    },
    "flac": {
        "mp3": lambda path: convert_audio_format(path, "mp3"),
        "wav": lambda path: convert_audio_format(path, "wav"),
        "aac": lambda path: convert_audio_format(path, "aac"),
        "ogg": lambda path: convert_audio_format(path, "ogg"),
        "opus": lambda path: convert_audio_format(path, "opus"),
        "ac3": lambda path: convert_audio_format(path, "ac3"),
        "m4a": lambda path: convert_audio_format(path, "m4a"),
        "alac": lambda path: convert_audio_format(path, "alac")
    },
    "opus": {
        "mp3": lambda path: convert_audio_format(path, "mp3"),
        "wav": lambda path: convert_audio_format(path, "wav"),
        "aac": lambda path: convert_audio_format(path, "aac"),
        "ogg": lambda path: convert_audio_format(path, "ogg"),
        "flac": lambda path: convert_audio_format(path, "flac"),
        "ac3": lambda path: convert_audio_format(path, "ac3"),
        "m4a": lambda path: convert_audio_format(path, "m4a"),
        "alac": lambda path: convert_audio_format(path, "alac")
    },
    "ac3": {
        "mp3": lambda path: convert_audio_format(path, "mp3"),
        "wav": lambda path: convert_audio_format(path, "wav"),
        "aac": lambda path: convert_audio_format(path, "aac"),
        "ogg": lambda path: convert_audio_format(path, "ogg"),
        "flac": lambda path: convert_audio_format(path, "flac"),
        "opus": lambda path: convert_audio_format(path, "opus"),
        "m4a": lambda path: convert_audio_format(path, "m4a"),
        "alac": lambda path: convert_audio_format(path, "alac")
    },
    "m4a": {
        "mp3": lambda path: convert_audio_format(path, "mp3"),
        "wav": lambda path: convert_audio_format(path, "wav"),
        "aac": lambda path: convert_audio_format(path, "aac"),
        "ogg": lambda path: convert_audio_format(path, "ogg"),
        "flac": lambda path: convert_audio_format(path, "flac"),
        "opus": lambda path: convert_audio_format(path, "opus"),
        "ac3": lambda path: convert_audio_format(path, "ac3"),
        "alac": lambda path: convert_audio_format(path, "alac")
    },
    "wma": {
        "mp3": lambda path: convert_audio_format(path, "mp3"),
        "wav": lambda path: convert_audio_format(path, "wav"),
        "aac": lambda path: convert_audio_format(path, "aac"),
        "ogg": lambda path: convert_audio_format(path, "ogg"),
        "flac": lambda path: convert_audio_format(path, "flac"),
        "opus": lambda path: convert_audio_format(path, "opus"),
        "ac3": lambda path: convert_audio_format(path, "ac3"),
        "m4a": lambda path: convert_audio_format(path, "m4a"),
        "alac": lambda path: convert_audio_format(path, "alac")
    },
    "heic": {
        "png": lambda path: convert_image_format(path, "png"),
        "jpg": lambda path: convert_image_format(path, "jpg"),
        "jpeg": lambda path: convert_image_format(path, "jpeg"),
        "webp": lambda path: convert_image_format(path, "webp"),
        "gif": lambda path: convert_image_format(path, "gif"),
        "bmp": lambda path: convert_image_format(path, "bmp"),
        "tiff": lambda path: convert_image_format(path, "tiff"),
        "ico": lambda path: convert_image_format(path, "ico"),
        "compressed": compress_image
    }
}

# After imports, add this helper function
def is_server_environment():
    """Detect if running in server environment to avoid desktop-only methods"""
    import platform
    import os
    
    # Check for container environment variable (set in our Dockerfile)
    if os.environ.get("RUNNING_IN_CONTAINER", "").lower() in ("1", "true", "yes"):
        print("Detected Docker container environment")
        return True
    
    # Non-Windows operating systems are definitely servers for our purposes
    # since PowerPoint COM automation only works on Windows
    if platform.system() != "Windows":
        return True
    
    # Check for common server indicators on Windows
    
    # Check for web server environment variables
    web_env_vars = ["SERVER_SOFTWARE", "REMOTE_ADDR", "HTTP_HOST", "REQUEST_URI", 
                    "GATEWAY_INTERFACE", "REQUEST_METHOD"]
    if any(env in os.environ for env in web_env_vars):
        return True
        
    # Check for common cloud/server providers' environment variables
    cloud_env_vars = ["WEBSITE_SITE_NAME", "AZURE_FUNCTIONS", "AWS_LAMBDA",
                     "DYNO", "HEROKU_APP_ID", "GOOGLE_CLOUD_PROJECT"]
    if any(env in os.environ for env in cloud_env_vars):
        return True
    
    # Additional server check: if running as a Windows service
    try:
        import win32ts
        session_id = win32ts.WTSGetActiveConsoleSessionId()
        # No active console session typically means a service
        if session_id == 0xFFFFFFFF:
            return True
    except:
        pass
    
    # Check if running without desktop interaction
    try:
        # If we can't access desktop-related functions, likely a server
        import win32gui
        desk = win32gui.GetDesktopWindow()
        # If this succeeds, probably not a server
    except:
        return True
    
    # Default to assuming it's not a server
    return False

# Make the platform detection available globally
running_on_server = is_server_environment()
if running_on_server:
    print("Running in server environment - desktop conversion methods will be skipped")
else:
    print("Running in desktop environment - all conversion methods will be available")

def direct_ffmpeg_conversion(input_path, output_format, ffmpeg_executable, task_id=None):
    """
    Synchronous function for direct FFmpeg audio extraction and conversion with detailed error logging
    """
    try:
        print(f"\n[DEBUG] Starting direct FFmpeg conversion from {input_path} to {output_format}")
        print(f"[DEBUG] FFmpeg path: {ffmpeg_executable}")
        
        if task_id and task_id in conversion_progress:
            conversion_progress[task_id] = {"progress": 40, "status": f"Using FFmpeg directly for {output_format} conversion"}
        
        # Create output path
        base_output_path = os.path.join(CONVERTED_DIR, os.path.splitext(os.path.basename(input_path))[0])
        output_path = f"{base_output_path}.{output_format}"
        print(f"[DEBUG] Target output path: {output_path}")
        
        # Check if ffmpeg exists
        if not os.path.exists(ffmpeg_executable):
            error_msg = f"FFmpeg executable not found at: {ffmpeg_executable}"
            print(f"[ERROR] {error_msg}")
            raise FileNotFoundError(error_msg)
        
        # Build the command - check if it's a video or audio file
        file_ext = os.path.splitext(input_path)[1].lower()
        is_video = file_ext in ['.mp4', '.avi', '.mov', '.mkv', '.webm', '.flv', '.wmv']
        
        cmd = [ffmpeg_executable, "-y", "-i", input_path]
        
        # For videos, we need to add -vn to strip video
        if is_video:
            cmd.append("-vn")
        
        # Format-specific options
        if output_format.lower() == "mp3":
            cmd.extend(["-c:a", "libmp3lame", "-q:a", "2"])
            cmd.append(output_path)
        elif output_format.lower() == "wav":
            cmd.extend(["-c:a", "pcm_s16le"])
            cmd.append(output_path)
        elif output_format.lower() == "aac":
            cmd.extend(["-c:a", "aac", "-strict", "experimental"])
            cmd.append(output_path) 
        elif output_format.lower() == "ogg":
            cmd.extend(["-c:a", "libvorbis"])
            cmd.append(output_path)
        elif output_format.lower() == "flac":
            cmd.extend(["-c:a", "flac"])
            cmd.append(output_path)
        elif output_format.lower() == "opus":
            cmd.extend(["-c:a", "libopus"])
            cmd.append(output_path)
        elif output_format.lower() == "ac3":
            cmd.extend(["-c:a", "ac3"])
            cmd.append(output_path)
        elif output_format.lower() == "m4a":
            output_path = f"{base_output_path}.m4a"
            cmd.extend(["-c:a", "aac", "-strict", "experimental"])
            cmd.append(output_path)
        elif output_format.lower() == "alac":
            output_path = f"{base_output_path}.m4a"
            cmd.extend(["-c:a", "alac"])
            cmd.append(output_path)
        else:
            # Default to MP3
            output_path = f"{base_output_path}.mp3"
            cmd.extend(["-c:a", "libmp3lame", "-q:a", "2"])
            cmd.append(output_path)
        
        # Print full command for debugging
        cmd_str = " ".join(cmd)
        print(f"[DEBUG] Running command: {cmd_str}")
        
        if task_id and task_id in conversion_progress:
            conversion_progress[task_id] = {"progress": 60, "status": "Running FFmpeg command"}
        
        # Run the command
        print("[DEBUG] Starting subprocess...")
        result = subprocess.run(
            cmd,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True
        )
        
        # Check result
        print(f"[DEBUG] Command exit code: {result.returncode}")
        if result.stdout:
            print(f"[DEBUG] Command stdout: {result.stdout}")
        if result.stderr:
            print(f"[DEBUG] Command stderr: {result.stderr}")
        
        # Verify file exists
        if os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            print(f"[DEBUG] Output file exists: {output_path}, size: {file_size} bytes")
        else:
            print(f"[ERROR] Output file does not exist: {output_path}")
        
        # Handle errors
        if result.returncode != 0 or not os.path.exists(output_path) or os.path.getsize(output_path) == 0:
            error_msg = f"FFmpeg conversion failed: {result.stderr}"
            print(f"[ERROR] {error_msg}")
            
            # Try fallback to MP3
            print("[DEBUG] Attempting fallback to MP3...")
            fallback_output = f"{base_output_path}.mp3"
            fallback_cmd = [ffmpeg_executable, "-y", "-i", input_path]
            
            # Add -vn flag only for video files
            if is_video:
                fallback_cmd.append("-vn")
                
            fallback_cmd.extend(["-c:a", "libmp3lame", "-q:a", "2", fallback_output])
            
            print(f"[DEBUG] Fallback command: {' '.join(fallback_cmd)}")
            fallback_result = subprocess.run(
                fallback_cmd,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True
            )
            
            print(f"[DEBUG] Fallback exit code: {fallback_result.returncode}")
            if fallback_result.stderr:
                print(f"[DEBUG] Fallback stderr: {fallback_result.stderr}")
            
            if fallback_result.returncode == 0 and os.path.exists(fallback_output) and os.path.getsize(fallback_output) > 0:
                print(f"[DEBUG] Fallback successful, using: {fallback_output}")
                output_path = fallback_output
            else:
                print(f"[ERROR] Fallback also failed: {fallback_result.stderr}")
                if task_id and task_id in conversion_progress:
                    conversion_progress[task_id] = {"progress": 100, "status": f"Error: FFmpeg conversion failed"}
                raise Exception(f"Failed to convert audio: {result.stderr}\nFallback also failed: {fallback_result.stderr}")
        
        # Update progress
        if task_id and task_id in conversion_progress:
            conversion_progress[task_id] = {
                "progress": 100, 
                "status": "Audio extraction complete",
                "file_path": output_path,
                "file_name": os.path.basename(output_path)
            }
        
        print(f"[DEBUG] Conversion successful: {output_path}")
        return output_path
        
    except Exception as e:
        print(f"[ERROR] Exception in direct_ffmpeg_conversion: {str(e)}")
        traceback_str = traceback.format_exc()
        print(f"[ERROR] Traceback: {traceback_str}")
        
        if task_id and task_id in conversion_progress:
            conversion_progress[task_id] = {
                "progress": 100, 
                "status": f"Error: {str(e)}"
            }
        
        raise Exception(f"Audio extraction failed: {str(e)}")

async def convert_audio_format(audio_path, target_format, task_id=None):
    """Convert audio from one format to another using FFmpeg"""
    try:
        print(f"Starting audio format conversion: {audio_path} to {target_format}")
        
        # Update progress if task_id provided
        if task_id:
            conversion_progress[task_id] = {"progress": 10, "status": "Starting audio conversion"}
        
        # Create output path based on input filename
        output_path = os.path.join(CONVERTED_DIR, os.path.splitext(os.path.basename(audio_path))[0] + f".{target_format}")
        
        # Check if ffmpeg is available
        if not ffmpeg_available:
            error_msg = "FFmpeg is required for audio conversion but not available"
            print(error_msg)
            if task_id:
                conversion_progress[task_id] = {"progress": 100, "status": f"Error: {error_msg}"}
            raise Exception(error_msg)
        
        # Get FFmpeg path - try the detected one or the local one
        ffmpeg_executable = ffmpeg_path
        if not ffmpeg_executable or not os.path.exists(ffmpeg_executable):
            local_ffmpeg = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ffmpeg", "ffmpeg.exe")
            if os.path.exists(local_ffmpeg):
                ffmpeg_executable = local_ffmpeg
        
        if not ffmpeg_executable or not os.path.exists(ffmpeg_executable):
            # Try to get it from imageio
            try:
                import imageio_ffmpeg
                ffmpeg_executable = imageio_ffmpeg.get_ffmpeg_exe()
            except Exception:
                pass
        
        if not ffmpeg_executable or not os.path.exists(ffmpeg_executable):
            error_msg = "Cannot locate FFmpeg executable"
            print(error_msg)
            if task_id:
                conversion_progress[task_id] = {"progress": 100, "status": f"Error: {error_msg}"}
            raise Exception(error_msg)
            
        if task_id:
            conversion_progress[task_id] = {"progress": 30, "status": "Preparing audio conversion"}
        
        # Use our direct FFmpeg conversion function for all audio conversions
        return direct_ffmpeg_conversion(audio_path, target_format, ffmpeg_executable, task_id)
    
    except Exception as e:
        print(f"Error converting audio format: {str(e)}")
        traceback_info = traceback.format_exc()
        print(f"Traceback: {traceback_info}")
        
        if task_id:
            conversion_progress[task_id] = {"progress": 100, "status": f"Error: {str(e)}"}
        
        raise Exception(f"Audio format conversion failed: {str(e)}")

async def convert_heic_specifically(image_path, target_format):
    """Direct HEIC to standard format conversion function bypassing complex logic"""
    print(f"=== SPECIAL HEIC CONVERSION DEBUG ===")
    print(f"Converting HEIC: {image_path} to {target_format}")
    print(f"Input file exists: {os.path.exists(image_path)}")
    print(f"Input file size: {os.path.getsize(image_path) if os.path.exists(image_path) else 'N/A'}")
    
    # Create output path
    output_filename = f"{os.path.splitext(os.path.basename(image_path))[0]}_converted.{target_format}"
    output_path = os.path.join(CONVERTED_DIR, output_filename)
    
    # First try to get dimensions with PIL after registering the HEIF opener
    try:
        import pillow_heif
        pillow_heif.register_heif_opener()
        
        img = Image.open(image_path)
        width, height = img.size
        print(f"PIL detected dimensions: {width}x{height}")
        
        # Convert directly with PIL which will preserve dimensions
        if img.mode not in ['RGB', 'RGBA'] and target_format not in ['bmp', 'gif']:
            img = img.convert('RGB')
        
        # Save with optimal quality settings
        save_opts = {}
        if target_format.lower() in ['jpg', 'jpeg']:
            save_opts = {'quality': 95, 'optimize': True, 'subsampling': 0}
            format_name = 'JPEG'  # Always use JPEG, not JPG
        elif target_format.lower() == 'png':
            save_opts = {'optimize': True}
            format_name = 'PNG'
        elif target_format.lower() == 'webp':
            save_opts = {'quality': 95, 'method': 6}
            format_name = 'WEBP'
        elif target_format.lower() == 'tiff':
            save_opts = {'compression': 'tiff_lzw'}
            format_name = 'TIFF'
        elif target_format.lower() == 'bmp':
            save_opts = {}
            format_name = 'BMP'
        elif target_format.lower() == 'gif':
            save_opts = {'optimize': True}
            format_name = 'GIF'
        elif target_format.lower() == 'pdf':
            # Handle PDF conversion directly
            print("Converting to PDF using direct approach")
            try:
                # Make sure image is in RGB mode for PDF
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                
                # Create a PDF directly from the image
                from reportlab.pdfgen import canvas
                from reportlab.lib.pagesizes import letter
                
                # Calculate dimensions to fit page while preserving aspect ratio
                width, height = letter
                img_width, img_height = img.size
                aspect = img_height / float(img_width)
                
                # Calculate dimensions to fit page
                display_width = width * 0.9  # 90% of page width
                display_height = display_width * aspect
                
                # Center on page
                x = (width - display_width) / 2
                y = (height - display_height) / 2
                
                # Save as temporary PNG first
                temp_png = os.path.join(TEMP_DIR, f"{os.path.splitext(os.path.basename(image_path))[0]}_temp.png")
                os.makedirs(TEMP_DIR, exist_ok=True)
                img.save(temp_png, format='PNG', optimize=True)
                
                # Create PDF
                c = canvas.Canvas(output_path, pagesize=letter)
                c.drawImage(temp_png, x, y, width=display_width, height=display_height)
                c.save()
                
                # Clean up temp file
                if os.path.exists(temp_png):
                    os.remove(temp_png)
                
                if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                    print(f"PDF created successfully via direct approach: {output_path}")
                    return output_path
                else:
                    print("PDF was not created successfully")
            except Exception as e:
                print(f"Error in direct PDF conversion: {str(e)}")
                traceback.print_exc()
            
            # If we get here, PDF conversion failed
            print("Direct PDF conversion failed")
            return None
        else:
            # Unknown format
            save_opts = {}
        
        # Save the image if not PDF
        if target_format.lower() != 'pdf':
            img.save(output_path, format=format_name, **save_opts)
        
        if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
            print(f"PIL conversion successful: {output_path}")
            return output_path
        else:
            print("PIL conversion failed, trying FFmpeg...")
    except Exception as e:
        print(f"PIL conversion failed: {str(e)}")
        width = height = None
    
    # If PIL fails or we couldn't get dimensions, try FFmpeg
    try:
        # Find FFmpeg
        ffmpeg_path = None
        try:
            import imageio_ffmpeg
            ffmpeg_path = imageio_ffmpeg.get_ffmpeg_exe()
        except:
            pass
        
        # Check for local FFmpeg
        if not ffmpeg_path:
            local_ffmpeg = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ffmpeg", "ffmpeg.exe")
            if os.path.exists(local_ffmpeg):
                ffmpeg_path = local_ffmpeg
        
        # Use FFmpeg if available
        if ffmpeg_path and os.path.exists(ffmpeg_path):
            print(f"Using FFmpeg for HEIC conversion: {ffmpeg_path}")
            
            # If we already have dimensions from PIL, use them
            if width and height:
                cmd = [
                    ffmpeg_path, 
                    "-y", 
                    "-i", image_path,
                    "-vf", f"scale={width}:{height}",
                    output_path
                ]
            else:
                # First, extract the original dimensions
                info_cmd = [
                    ffmpeg_path,
                    "-i", image_path
                ]
                
                info_result = subprocess.run(info_cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
                print(f"FFmpeg info stderr: {info_result.stderr}")
                
                # Extract dimensions using various patterns
                dimensions_match = None
                
                # Try specific pattern for tile grid
                tile_match = re.search(r'Stream group.+?(\d{3,})x(\d{3,})', info_result.stderr)
                if tile_match:
                    dimensions_match = tile_match
                
                # Try general pattern for video streams
                if not dimensions_match:
                    video_match = re.search(r'Video:.+?(\d{3,})x(\d{3,})', info_result.stderr)
                    if video_match:
                        dimensions_match = video_match
                
                # Fall back to any 3+ digit x 3+ digit pattern
                if not dimensions_match:
                    any_match = re.search(r'(\d{3,})x(\d{3,})', info_result.stderr)
                    if any_match:
                        dimensions_match = any_match
                
                if dimensions_match:
                    width = dimensions_match.group(1)
                    height = dimensions_match.group(2)
                    
                    # Verify dimensions
                    try:
                        width_int = int(width)
                        height_int = int(height)
                        
                        if width_int < 100 or height_int < 100:
                            print(f"Suspicious dimensions: {width}x{height}, using basic command")
                            cmd = [
                                ffmpeg_path, 
                                "-y", 
                                "-i", image_path,
                                output_path
                            ]
                        else:
                            print(f"FFmpeg detected dimensions: {width}x{height}")
                            cmd = [
                                ffmpeg_path, 
                                "-y", 
                                "-i", image_path,
                                "-vf", f"scale={width}:{height}",
                                output_path
                            ]
                    except ValueError:
                        print(f"Invalid dimensions: {width}x{height}, using basic command")
                        cmd = [
                            ffmpeg_path, 
                            "-y", 
                            "-i", image_path,
                            output_path
                        ]
                else:
                    print("Couldn't detect dimensions, using basic command")
                    cmd = [
                        ffmpeg_path, 
                        "-y", 
                        "-i", image_path,
                        output_path
                    ]
            
            print(f"FFmpeg command: {' '.join(cmd)}")
            result = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
            print(f"FFmpeg stdout: {result.stdout}")
            print(f"FFmpeg stderr: {result.stderr}")
            print(f"FFmpeg exit code: {result.returncode}")
            
            if result.returncode == 0 and os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                print(f"FFmpeg HEIC conversion successful: {output_path}")
                return output_path
            else:
                print(f"FFmpeg conversion failed, trying HeifFile approach")
    except Exception as e:
        print(f"FFmpeg attempt failed: {str(e)}")
    
    # If both PIL and FFmpeg fail, try the explicit HeifFile approach
    try:
        import pillow_heif
        print(f"Using pillow-heif version: {pillow_heif.__version__}")
        
        # Try with explicit HeifFile approach
        try:
            print("Trying explicit HeifFile approach")
            heif_file = pillow_heif.open_heif(image_path)
            print(f"HEIF image info - Mode: {heif_file.mode}, Size: {heif_file.size}")
            
            # Create a PIL image from the HEIF data
            img = Image.frombytes(
                heif_file.mode,
                heif_file.size,
                heif_file.data,
                "raw",
            )
            
            print(f"Successfully created PIL image from HeifFile. Mode: {img.mode}, Size: {img.size}")
            
            # Convert mode if needed
            if img.mode not in ['RGB', 'RGBA'] and target_format not in ['bmp', 'gif']:
                print(f"Converting image mode from {img.mode} to RGB")
                img = img.convert('RGB')
            
            # Format-specific options for best quality
            save_opts = {}
            if target_format.lower() in ['jpg', 'jpeg']:
                save_opts = {'quality': 95, 'optimize': True, 'subsampling': 0}  # Highest quality JPEG
                format_name = 'JPEG'  # Always use JPEG, not JPG
            elif target_format.lower() == 'png':
                save_opts = {'optimize': True, 'compress_level': 1}  # Best quality PNG
                format_name = 'PNG'
            elif target_format.lower() == 'webp':
                save_opts = {'quality': 95, 'method': 6, 'lossless': True}  # High quality WebP
                format_name = 'WEBP'
            elif target_format.lower() == 'tiff':
                save_opts = {'compression': 'tiff_lzw'}  # Good compression for TIFF
                format_name = 'TIFF'
            elif target_format.lower() == 'bmp':
                save_opts = {}
                format_name = 'BMP'
            elif target_format.lower() == 'gif':
                save_opts = {'optimize': True}
                format_name = 'GIF'
            elif target_format.lower() == 'pdf':
                # Handle PDF conversion separately - for PDF we'll skip this approach
                # and let the FFmpeg method handle it
                print("HeifFile approach is skipping PDF format - will try FFmpeg instead")
                return None  # Skip to next method
            else:
                # Unknown format
                save_opts = {}
                format_name = target_format.upper()
            
            # Save the image if not PDF
            if target_format.lower() != 'pdf':
                img.save(output_path, format=format_name, **save_opts)
            
            # Verify the output
            if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                print(f"HeifFile approach successful: {output_path}")
                return output_path
            else:
                print(f"HeifFile approach failed to create output file")
                raise Exception("Failed to create output file with HeifFile approach")
        
        except Exception as e2:
            print(f"Error with HeifFile approach: {str(e2)}")
            raise Exception(f"HEIC conversion failed: All approaches failed")
    
    except Exception as e:
        print(f"All HEIC conversion approaches failed: {str(e)}")
        traceback.print_exc()
        raise Exception(f"HEIC conversion failed: {str(e)}")

# Update conversion_functions for HEIC
conversion_functions["heic"] = {
    "png": lambda path: convert_heic_specifically(path, "png"),
    "jpg": lambda path: convert_heic_specifically(path, "jpg"),
    "jpeg": lambda path: convert_heic_specifically(path, "jpeg"),
    "webp": lambda path: convert_heic_specifically(path, "webp"),
    "gif": lambda path: convert_heic_specifically(path, "gif"),
    "bmp": lambda path: convert_heic_specifically(path, "bmp"),
    "tiff": lambda path: convert_heic_specifically(path, "tiff"),
    "ico": lambda path: convert_heic_specifically(path, "ico"),
    "compressed": compress_image
}

async def direct_heic_to_pdf(image_path, output_path):
    """Direct conversion from HEIC to PDF bypassing convert_image_to_pdf"""
    print(f"Direct HEIC to PDF conversion: {image_path} -> {output_path}")
    try:
        # Register HEIF opener and open the image
        import pillow_heif
        pillow_heif.register_heif_opener()
        
        img = Image.open(image_path)
        print(f"Opened HEIC image. Size: {img.size}, Mode: {img.mode}")
        
        # Convert to RGB for PDF
        if img.mode != 'RGB':
            img = img.convert('RGB')
            print("Converted image to RGB mode")
        
        # Save as temporary PNG
        temp_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "temp")
        os.makedirs(temp_dir, exist_ok=True)
        temp_png = os.path.join(temp_dir, f"{os.path.splitext(os.path.basename(image_path))[0]}_temp.png")
        
        print(f"Saving temporary PNG: {temp_png}")
        img.save(temp_png, format='PNG', optimize=True)
        
        if not os.path.exists(temp_png) or os.path.getsize(temp_png) == 0:
            print("Failed to create temporary PNG file")
            return None
        
        # Create PDF using reportlab
        try:
            from reportlab.pdfgen import canvas
            from reportlab.lib.pagesizes import letter
            
            print(f"Creating PDF with ReportLab: {output_path}")
            
            # Calculate dimensions to fit page while preserving aspect ratio
            width, height = letter
            img_width, img_height = img.size
            aspect = img_height / float(img_width)
            
            # Calculate dimensions to fit page
            display_width = width * 0.9  # 90% of page width
            display_height = display_width * aspect
            
            # Center on page
            x = (width - display_width) / 2
            y = (height - display_height) / 2
            
            # Create PDF
            c = canvas.Canvas(output_path, pagesize=letter)
            c.drawImage(temp_png, x, y, width=display_width, height=display_height)
            c.save()
            
            print(f"PDF created successfully: {output_path}")
        except Exception as e:
            print(f"ReportLab PDF creation failed: {str(e)}")
            
            # Fallback to img2pdf if reportlab fails
            try:
                import img2pdf
                print(f"Trying img2pdf for PDF creation: {output_path}")
                
                with open(output_path, "wb") as f:
                    f.write(img2pdf.convert(temp_png))
                
                print(f"img2pdf PDF creation successful: {output_path}")
            except Exception as e2:
                print(f"img2pdf also failed: {str(e2)}")
                return None
        
        # Clean up temp file
        if os.path.exists(temp_png):
            os.remove(temp_png)
            print(f"Removed temporary PNG file: {temp_png}")
        
        # Check if PDF was created successfully
        if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
            print(f"PDF file created successfully: {output_path}")
            return output_path
        else:
            print(f"PDF file not created or empty: {output_path}")
            return None
    
    except Exception as e:
        print(f"Error in direct_heic_to_pdf: {str(e)}")
        traceback.print_exc()
        return None

# Update the conversion functions dictionary for HEIC to PDF
conversion_functions["heic"]["pdf"] = lambda path: direct_heic_to_pdf(path, os.path.join(CONVERTED_DIR, f"{os.path.splitext(os.path.basename(path))[0]}_converted.pdf"))

if __name__ == "__main__":
    import argparse
    import uvicorn
    from socket import gethostname, gethostbyname
    
    # Parse command line arguments
    parser = argparse.ArgumentParser(description="FileAlchemy - Your Ultimate Document Conversion Solution")
    parser.add_argument('--host', type=str, default="0.0.0.0", 
                        help="Host address to bind the server (default: 0.0.0.0 for Railway)")
    parser.add_argument('--port', type=int, default=PORT, help=f"Port to bind the server (default: {PORT} from Railway)")
    parser.add_argument('--share', action='store_true', help="Create a public URL using ngrok (requires ngrok installation)")
    args = parser.parse_args()
    
    # Print server information
    print(f"\n{'='*50}")
    print(f"FileAlchemy Server Starting")
    print(f"{'='*50}")
    
    # Check conversion readiness
    print("\nChecking conversion capabilities:")
    available_features = []
    limitations = []
    
    # Check PDF libraries
    try:
        import fitz
        available_features.append("PDF processing (PyMuPDF)")
    except ImportError:
        limitations.append("PDF processing limited (PyMuPDF not found)")
    
    # Check image libraries
    try:
        import PIL
        available_features.append("Image processing (PIL/Pillow)")
    except ImportError:
        limitations.append("Image processing limited (PIL/Pillow not found)")
    
    # Check OCR
    try:
        import easyocr
        available_features.append("OCR capabilities (EasyOCR)")
    except ImportError:
        limitations.append("OCR not available (EasyOCR not found)")
    
    # Check DOCX
    try:
        from docx import Document
        available_features.append("Word document processing (python-docx)")
    except ImportError:
        limitations.append("Word document processing limited (python-docx not found)")
    
    # Check video
    if ffmpeg_available:
        available_features.append("Video conversion (FFmpeg)")
    else:
        limitations.append("Video conversion limited (FFmpeg not found, will create GIFs instead)")
    
    # Check HEIC support
    try:
        import pillow_heif
        print(f"\n✅ HEIC support detected. pillow-heif version: {pillow_heif.__version__}")
        available_features.append("HEIC image support (pillow-heif)")
        
        # Test HEIC conversion capability
        try:
            test_img = Image.new('RGB', (10, 10), color='red')
            from pillow_heif import from_pillow
            heif_file = from_pillow(test_img)
            test_heic_path = os.path.join(TEMP_DIR, "test_heic.heic")
            os.makedirs(TEMP_DIR, exist_ok=True)
            heif_file.save(test_heic_path)
            if os.path.exists(test_heic_path) and os.path.getsize(test_heic_path) > 0:
                print(f"✅ HEIC conversion test successful. Test file created at: {test_heic_path}")
            else:
                print("⚠️ HEIC test file was not created successfully")
                limitations.append("HEIC conversion may not work properly")
        except Exception as e:
            print(f"⚠️ HEIC conversion test failed: {str(e)}")
            limitations.append(f"HEIC conversion may not work properly: {str(e)}")
    except ImportError:
        print("\n⚠️ HEIC support not available (pillow-heif not found)")
        limitations.append("HEIC image support limited (pillow-heif not found)")
    except Exception as e:
        print(f"\n⚠️ HEIC support error: {str(e)}")
        limitations.append(f"HEIC image support error: {str(e)}")
    
    # Print features and limitations
    print("\nAvailable features:")
    for feature in available_features:
        print(f"✅ {feature}")
    
    if limitations:
        print("\nLimitations:")
        for limitation in limitations:
            print(f"⚠️ {limitation}")
    
    # Print access information
    local_ip = "127.0.0.1"
    try:
        # Try to get local network IP for LAN access
        hostname = gethostname()
        local_ip = gethostbyname(hostname)
    except:
        pass
        
    print(f"\nAccess information:")
    print(f"- Local URL: http://{args.host}:{args.port}")
    
    if args.host == "0.0.0.0":
        print(f"- Network URL: http://{local_ip}:{args.port} (for devices on same network)")
    
    # Start ngrok tunnel if requested
    if args.share:
        # First check if pyngrok is installed without importing it
        pyngrok_available = False
        try:
            import importlib.util
            pyngrok_spec = importlib.util.find_spec("pyngrok")
            pyngrok_available = pyngrok_spec is not None
        except ImportError:
            pyngrok_available = False
        
        if not pyngrok_available:
            print("\n❌ Failed to create public URL: pyngrok not installed")
            print("To enable public sharing, install with: pip install pyngrok")
        else:
            try:
                from pyngrok import ngrok
                
                # Open a tunnel to the uvicorn server
                public_url = ngrok.connect(args.port).public_url
                print(f"- Public URL: {public_url} (for sharing outside your network)")
                print("\n⚠️ Warning: Your files will be accessible to anyone with this URL!")
            except Exception as e:
                print(f"\n❌ Failed to create public URL: {str(e)}")
                print("Make sure ngrok is installed and properly configured")
    
    print(f"\n{'='*50}")    

    # Start the server
    print(f"\nStarting FileAlchemy server...")
    print(f"Environment: {RAILWAY_ENVIRONMENT}")
    print(f"Host: {args.host}")
    print(f"Port: {args.port}")
    
    # Run the server
    uvicorn.run(
        "server:app",
        host=args.host,
        port=args.port,
        reload=False,  # Disable reload in production
        access_log=True,
        log_level="info"
    )